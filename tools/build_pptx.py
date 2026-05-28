"""
Génère TrafficSimulation_Presentation.pptx depuis docs/.
Couvre : état de l'art, comportements, expériences, architecture.
Inclut placeholders image/vidéo + notes (script complet).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- Palette ---
C_BG     = RGBColor(0x10, 0x14, 0x22)
C_FG     = RGBColor(0xF2, 0xF4, 0xFA)
C_ACC    = RGBColor(0x3D, 0x9D, 0xF3)
C_ACC2   = RGBColor(0xFF, 0xB7, 0x3D)
C_OK     = RGBColor(0x4C, 0xD9, 0x7F)
C_BAD    = RGBColor(0xF1, 0x66, 0x66)
C_DIM    = RGBColor(0x8A, 0x90, 0xA8)
C_PH     = RGBColor(0x24, 0x2A, 0x3C)
C_PH_BD  = RGBColor(0x3D, 0x9D, 0xF3)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ----------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------
def add_bg(slide, color=C_BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    return bg

def add_accent_bar(slide, color=C_ACC):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SH)
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    return bar

def add_text(slide, x, y, w, h, text, *, size=14, bold=False,
             color=C_FG, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def add_title(slide, title, subtitle=None, *, chapter=None):
    add_accent_bar(slide)
    if chapter:
        add_text(slide, Inches(0.4), Inches(0.25), Inches(12), Inches(0.32),
                 chapter, size=11, bold=True, color=C_ACC)
    add_text(slide, Inches(0.4), Inches(0.55), Inches(12.5), Inches(0.85),
             title, size=30, bold=True, color=C_FG)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(1.32), Inches(12.5), Inches(0.45),
                 subtitle, size=15, color=C_DIM)
    # separator
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.4), Inches(1.85),
                                 Inches(12.5), Emu(20000))
    sep.line.fill.background()
    sep.fill.solid(); sep.fill.fore_color.rgb = C_ACC

def add_bullets(slide, x, y, w, h, items, *, size=14, color=C_FG, bullet="•"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {it}"
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb

def add_card(slide, x, y, w, h, title, body, *, accent=C_ACC):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.line.color.rgb = accent
    card.line.width = Pt(1.0)
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x18, 0x1E, 0x30)
    card.shadow.inherit = False
    # title
    add_text(slide, x + Inches(0.15), y + Inches(0.1),
             w - Inches(0.3), Inches(0.35),
             title, size=13, bold=True, color=accent)
    # body
    add_text(slide, x + Inches(0.15), y + Inches(0.5),
             w - Inches(0.3), h - Inches(0.55),
             body, size=11, color=C_FG)

def add_placeholder(slide, x, y, w, h, label="IMAGE / VIDÉO"):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.line.color.rgb = C_PH_BD
    box.line.width = Pt(1.2)
    box.line.dash_style = 7  # dash
    box.fill.solid(); box.fill.fore_color.rgb = C_PH
    add_text(slide, x, y + h/2 - Inches(0.25), w, Inches(0.5),
             f"[ {label} ]", size=13, bold=True,
             color=C_PH_BD, align=PP_ALIGN.CENTER)
    return box

def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_footer(slide, page, total, chapter=""):
    add_text(slide, Inches(0.4), SH - Inches(0.35),
             Inches(6), Inches(0.3),
             f"TrafficSimulation MAS · {chapter}", size=9, color=C_DIM)
    add_text(slide, SW - Inches(1.2), SH - Inches(0.35),
             Inches(0.8), Inches(0.3),
             f"{page} / {total}", size=9, color=C_DIM,
             align=PP_ALIGN.RIGHT)

# ----------------------------------------------------------------------
# Slide builders
# ----------------------------------------------------------------------
SLIDES = []

def slide_cover(title, subtitle, author, date):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    # gradient-ish accent
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              0, Inches(2.8), SW, Inches(0.06))
    band.line.fill.background()
    band.fill.solid(); band.fill.fore_color.rgb = C_ACC
    add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
             "PRÉSENTATION DE RECHERCHE", size=14, bold=True, color=C_ACC2)
    add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(1.4),
             title, size=42, bold=True, color=C_FG)
    add_text(s, Inches(0.6), Inches(3.1), Inches(12), Inches(0.7),
             subtitle, size=18, color=C_DIM)
    add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.35),
             f"{author}   ·   {date}", size=12, color=C_DIM)
    SLIDES.append((s, "Couverture"))
    return s

def slide_section(num, title, subtitle):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.6),
             f"PARTIE {num}", size=18, bold=True, color=C_ACC2)
    add_text(s, Inches(0.6), Inches(2.6), Inches(12), Inches(1.4),
             title, size=44, bold=True, color=C_FG)
    add_text(s, Inches(0.6), Inches(4.0), Inches(12), Inches(0.6),
             subtitle, size=16, color=C_DIM)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.6), Inches(4.85), Inches(2.5), Emu(40000))
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = C_ACC
    SLIDES.append((s, f"Partie {num}"))
    return s

def slide_content(chapter, title, subtitle=None):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_title(s, title, subtitle, chapter=chapter)
    SLIDES.append((s, chapter))
    return s

# ======================================================================
# CONTENT
# ======================================================================
TITLE = "TrafficSimulation MAS"
SUB   = ("Simulation multi-agents du trafic urbain : "
         "état de l'art, implémentation, expériences et architecture")
AUTH  = "Auteur : David François"
DATE  = "2026"

# ---------- 0. COVER + PLAN ----------
slide_cover(TITLE, SUB, AUTH, DATE)
set_notes(SLIDES[-1][0],
"Bonjour. Je vais vous présenter mon projet de simulation multi-agents du trafic urbain. "
"L'objectif est de transformer une preuve de concept logicielle en une preuve de "
"technologie scientifique : un banc d'essai capable de comparer plusieurs algorithmes "
"de coordination aux intersections, classiques et issus de la recherche récente. "
"Je vais d'abord poser l'état de l'art, puis présenter l'implémentation et tous "
"les comportements qui fonctionnent, ensuite les expériences menées et les résultats "
"comparés à la littérature, et enfin l'architecture logicielle complète.")

s = prs.slides.add_slide(BLANK); add_bg(s)
add_title(s, "Plan de la présentation", "Cinq parties + conclusion",
          chapter="AGENDA")
items = [
 ("1.  État de l'art scientifique",
  "IDM, AIM, P2P VanMiddlesworth, ORCA, Virtual Platooning, MAPF, MARL"),
 ("2.  Développement & comportements implémentés",
  "10 régulations + 15 fiches comportement + FSM décisionnelle"),
 ("3.  Observations & expériences",
  "Monte-Carlo, métriques, comportements émergents, concordance littérature"),
 ("4.  Architecture logicielle",
  "Pipeline 2-phases, Strategy, diagrammes de classes, threading"),
 ("5.  Limitations & travaux futurs",
  "Blocages résiduels grosse simulation, MOBIL multi-voies, MARL, LLM-agents"),
]
y = Inches(2.1)
for tag, sub in items:
    add_text(s, Inches(0.7), y, Inches(11.5), Inches(0.38),
             tag, size=18, bold=True, color=C_ACC)
    add_text(s, Inches(0.7), y + Inches(0.42), Inches(11.5), Inches(0.4),
             sub, size=13, color=C_DIM)
    y += Inches(0.95)
SLIDES.append((s, "Agenda"))
set_notes(s,
"La présentation suit cinq mouvements. Je commencerai par établir le cadre "
"scientifique : pourquoi un simulateur de trafic moderne ne peut plus se contenter "
"de feux et de panneaux STOP, mais doit intégrer des agents autonomes capables de "
"négocier. Ensuite, je vais montrer concrètement ce qui tourne dans mon application : "
"dix algorithmes de coordination différents, quinze comportements de conduite, "
"tous interchangeables à chaud. Puis je présenterai le banc d'essai Monte-Carlo et "
"les résultats. Enfin, je détaillerai l'architecture logicielle qui rend tout cela "
"propre et extensible. Je terminerai par une discussion honnête des limitations.")

# ======================================================================
# PART 1 — ÉTAT DE L'ART
# ======================================================================
slide_section(1, "État de l'art",
"Du PoC logiciel à la preuve de technologie scientifique")
set_notes(SLIDES[-1][0],
"Cette première partie est purement scientifique. Je vais résumer les articles "
"fondateurs qui structurent le domaine de la coordination multi-agents aux "
"intersections, et qui ont directement inspiré l'architecture de mon simulateur.")

# 1.1 Paradigme
s = slide_content("ÉTAT DE L'ART · 1.1",
                  "Changement de paradigme : PoC → PoT",
                  "Pourquoi une simulation \"avec des stops\" ne suffit plus")
add_bullets(s, Inches(0.6), Inches(2.1), Inches(7.6), Inches(4.5), [
 "Un PoC logiciel : véhicules cinématiques + règles statiques (stop, yield, feux).",
 "Une PoT scientifique : agents autonomes, perception locale, protocoles de négociation.",
 "Problématique : comment des règles microscopiques produisent des phénomènes macroscopiques mesurables (débit, délai, sécurité) ?",
 "Méthode : injecter de la variabilité stochastique, mesurer, comparer aux prédictions théoriques.",
 "Domaines mobilisés : SMA, ingénierie du trafic, théorie des jeux, apprentissage par renforcement.",
], size=14)
add_placeholder(s, Inches(8.5), Inches(2.1), Inches(4.4), Inches(3.0),
                "Schéma : PoC vs PoT")
add_text(s, Inches(8.5), Inches(5.25), Inches(4.4), Inches(1.6),
"Référence : VanMiddlesworth, Dresner & Stone, AAMAS 2008.\n"
"« Replacing the Stop Sign ».",
size=11, color=C_DIM)
set_notes(s,
"On part d'une distinction fondamentale faite par Dresner et Stone : un simulateur de "
"trafic avec des règles fixes — STOP, cédez le passage, feux — est une preuve de "
"concept logicielle. Pour faire de la recherche, il faut passer à la preuve de "
"technologie : chaque véhicule devient un agent autonome qui perçoit, raisonne et "
"négocie. La question scientifique devient alors : quels protocoles microscopiques de "
"coordination produisent les meilleures métriques macroscopiques ? C'est précisément "
"cette question que mon banc d'essai cherche à instrumenter.")

# 1.2 IDM
s = slide_content("ÉTAT DE L'ART · 1.2",
                  "Intelligent Driver Model (Treiber & Helbing, 2000)",
                  "Le standard incontesté de la dynamique longitudinale")
add_card(s, Inches(0.6), Inches(2.1), Inches(6.1), Inches(2.6),
"Équation d'accélération",
"a = a_max · [ 1 − (v / v₀)^δ − (s*(v,Δv) / s)² ]\n"
"s*(v,Δv) = s₀ + max(0, v·T + v·Δv / (2·√(a·b)))\n\n"
"Paramètres : v₀ (vitesse libre), T (headway), s₀ (jam distance),\n"
"a (acc. max confortable), b (décél. confortable), δ (exposant).",
accent=C_ACC)
add_card(s, Inches(6.85), Inches(2.1), Inches(6.05), Inches(2.6),
"Propriétés clés",
"• Continuité et différentiabilité (jerk fini).\n"
"• String stability : amortit les ondes de freinage fantômes.\n"
"• Transition fluide free-flow → car-following.\n"
"• Ligne d'arrêt = leader virtuel à v=0 → freinage naturel.\n"
"• Stochastique possible (T gaussien) pour hétérogénéité.",
accent=C_ACC2)
add_placeholder(s, Inches(0.6), Inches(4.85), Inches(12.3), Inches(2.1),
                "Courbe : profil de vitesse IDM (free-flow → arrêt à ligne virtuelle)")
set_notes(s,
"L'Intelligent Driver Model est le contrôleur longitudinal de chaque véhicule. "
"C'est un modèle continu, différentiable, qui produit des décélérations "
"asymptotiques douces. Sa force pour notre projet : il transforme n'importe "
"quelle contrainte — un véhicule réel, une ligne d'arrêt, un meneur virtuel — "
"en obstacle perçu, et calcule l'accélération appropriée. Aucune logique de "
"freinage explicite à coder : tout passe par le modèle. La propriété de string "
"stability garantit qu'un coup de frein du premier véhicule ne provoque pas "
"un arrêt complet plusieurs centaines de mètres en arrière.")

# 1.3 Taxonomy paradigmes
s = slide_content("ÉTAT DE L'ART · 1.3",
                  "Trois paradigmes de gestion d'intersection",
                  "Centralisé / décentralisé pair-à-pair / géométrique continu")
cards = [
 ("AIM (centralisé)",
  "Dresner & Stone, JAIR 2008. Un Intersection Manager arbitre des réservations "
  "FCFS sur une grille spatio-temporelle de tuiles. Quasi-zéro attente à densité "
  "modérée, mais point de défaillance unique + infrastructure dédiée.", C_ACC),
 ("P2P (décentralisé)",
  "VanMiddlesworth, Dresner & Stone, AAMAS 2008. Négociation broadcast via VANET. "
  "Hiérarchie de dominance locale, aucun manager. Optimal à faible densité ; "
  "dégradation au-delà de ~0.35 véh/s/voie vers comportement four-way stop.", C_ACC2),
 ("ORCA (géométrique continu)",
  "Van den Berg et al. 2011. Réciprocité de demi-plans de vélocité. Initialement "
  "robotique / foule, adapté ici à des voies curvilignes. Aucune file d'attente : "
  "tout est modulation continue de la vitesse cible.", C_OK),
]
y = Inches(2.1)
for t, b, c in cards:
    add_card(s, Inches(0.6), y, Inches(12.3), Inches(1.45), t, b, accent=c)
    y += Inches(1.6)
set_notes(s,
"L'état de l'art se structure en trois familles. AIM, c'est le paradigme centralisé : "
"un manager omniscient distribue des créneaux spatio-temporels. Performant mais "
"vulnérable et coûteux. P2P, c'est le pari de la décentralisation totale : chaque "
"agent négocie avec ses voisins via VANET. Aucune infrastructure, mais une "
"saturation au-delà d'un seuil de densité. ORCA, c'est l'approche purement "
"géométrique, sans signaling : les agents se croisent en modulant continûment "
"leur vitesse, comme une foule ou un essaim de drones. Mon simulateur implémente "
"les trois familles, plus les régulations classiques, pour comparer.")

# 1.4 AIM
s = slide_content("ÉTAT DE L'ART · 1.4",
                  "AIM — Autonomous Intersection Management",
                  "Dresner & Stone (JAIR 2008)")
add_bullets(s, Inches(0.6), Inches(2.1), Inches(7.4), Inches(4.5), [
 "Manager centralisé : reçoit (v, a, dimensions, ETA, trajectoire) de chaque agent.",
 "Discrétise le carrefour en grille de tuiles spatio-temporelles.",
 "Politique FCFS : simule la trajectoire requise, accepte si aucun chevauchement.",
 "Refus → l'agent ralentit via IDM, re-soumet une nouvelle fenêtre.",
 "Empiriquement : attentes ≈ 0 s à densité modérée, vitesses élevées maintenues.",
 "Extension : priorité urgence (ambulance) forçant les autres à avorter.",
 "Limites : point de défaillance unique, infrastructure carrefour coûteuse.",
], size=13)
add_placeholder(s, Inches(8.2), Inches(2.1), Inches(4.7), Inches(3.8),
                "Schéma : grille de réservation tile×time")
add_text(s, Inches(8.2), Inches(6.05), Inches(4.7), Inches(0.9),
"Réf. : Dresner K., Stone P., « A Multiagent Approach to Autonomous "
"Intersection Management », JAIR 2008.", size=10, color=C_DIM)
set_notes(s,
"AIM, c'est l'approche centralisée formalisée par Dresner et Stone. Chaque véhicule "
"envoie sa demande au manager, qui simule mentalement la trajectoire dans une grille "
"spatio-temporelle. Si toutes les cellules sont libres dans la fenêtre demandée, la "
"réservation est accordée. Le véhicule s'engage alors à respecter exactement le "
"profil cinématique annoncé. Sinon, il ralentit et re-soumet. Les résultats publiés "
"sont spectaculaires : à densité moyenne, on observe quasi-zéro temps d'attente. "
"Les véhicules se croisent à pleine vitesse en se manquant de quelques mètres. "
"Mais cela suppose une infrastructure dédiée et un manager fiable — d'où l'intérêt "
"des approches décentralisées qui suivent.")

# 1.5 P2P protocol
s = slide_content("ÉTAT DE L'ART · 1.5",
                  "P2P : protocole pair-à-pair (VanMiddlesworth 2008)",
                  "Aucun manager, broadcast asynchrone via VANET")
add_text(s, Inches(0.6), Inches(2.1), Inches(12.3), Inches(0.35),
"Machine à 3 états :", size=14, bold=True, color=C_ACC2)
states = [
 ("LURKING",  "À ~200 m. L'agent ÉCOUTE les Claims des pairs, construit la carte "
              "mentale du carrefour. Aucun envoi → évite la broadcast storm."),
 ("CLAIMING", "Diffuse un Claim {VIN, lane_index, intent, estimated_exit_time, "
              "is_stopped}. Compare aux Claims reçus selon la hiérarchie de dominance."),
 ("TRAVERSAL","Passage actif. Le Claim est retiré du bus dès la sortie. "
              "Si dominé pendant Claiming → Cancel + freinage IDM."),
]
y = Inches(2.55)
for n, b in states:
    add_card(s, Inches(0.6), y, Inches(12.3), Inches(1.05), n, b, accent=C_ACC)
    y += Inches(1.2)
add_placeholder(s, Inches(0.6), Inches(6.15), Inches(12.3), Inches(0.95),
                "Diagramme : transitions LURKING → CLAIMING → TRAVERSAL")
set_notes(s,
"Le protocole P2P de VanMiddlesworth s'articule autour d'une machine à trois "
"états. En LURKING, à environ 200 mètres du carrefour, le véhicule écoute en "
"silence. Cette latence intentionnelle évite la tempête de broadcast et garantit "
"qu'il ne fait pas de revendication sur des informations obsolètes. En CLAIMING, "
"plus près, il diffuse son Claim — son intention complète. À chaque tick, il "
"compare son propre Claim à ceux reçus selon une hiérarchie stricte. S'il est "
"dominé, il annule, freine via IDM. En TRAVERSAL, le passage est en cours, "
"on retire le Claim une fois sorti. Mon implémentation reprend exactement cette "
"FSM, avec en plus le verrou commit-to-pass pour éviter qu'un véhicule "
"engagé re-cède au milieu du carrefour.")

# 1.6 P2P hierarchy
s = slide_content("ÉTAT DE L'ART · 1.6",
                  "Hiérarchie de dominance P2P",
                  "Règles strictes pour éviter les deadlocks")
hier = [
 ("1. Temporelle",
  "Les deux en mouvement → plus petit estimated_exit_time gagne."),
 ("2. Spatiale (droite)",
  "Les deux à l'arrêt → véhicule géométriquement à droite gagne."),
 ("3. Manœuvre",
  "Symétrie → trajectoire rectiligne > trajectoire en virage."),
 ("4. VIN arbitraire",
  "Symétrie absolue → plus petit Vehicle ID gagne. Brise tout cycle."),
]
y = Inches(2.15)
for n, b in hier:
    add_card(s, Inches(0.6), y, Inches(12.3), Inches(0.95), n, b, accent=C_ACC2)
    y += Inches(1.10)
add_text(s, Inches(0.6), Inches(6.7), Inches(12.3), Inches(0.5),
"Seuil critique observé : ~0.35 véh/s/voie. Au-delà, P2P dégénère vers "
"comportement four-way stop (paradoxe de Braess).",
size=12, color=C_BAD, bold=True)
set_notes(s,
"La hiérarchie de dominance est ce qui rend P2P viable. C'est un arbre de "
"décision strict, identique chez tous les agents, qui garantit qu'au moins un "
"acteur cède dans chaque conflit. La règle ultime est le bris d'égalité par "
"VIN : le plus petit identifiant l'emporte. C'est trivialement déterministe, "
"donc reproductible en Monte-Carlo. J'ai repris exactement cette hiérarchie "
"dans ma P2PPolicy. La limitation théorique est bien documentée : au-delà "
"d'environ 0.35 véhicule par seconde par voie, la fréquence des conflits "
"sature la négociation et P2P s'effondre vers un comportement de four-way stop. "
"C'est précisément ce que mon banc d'essai cherche à reproduire empiriquement.")

# 1.7 ORCA
s = slide_content("ÉTAT DE L'ART · 1.7",
                  "ORCA — Optimal Reciprocal Collision Avoidance",
                  "Van den Berg et al. 2011")
add_bullets(s, Inches(0.6), Inches(2.1), Inches(7.6), Inches(4.5), [
 "Chaque agent = disque dans l'espace de vélocité 2D.",
 "Velocity Obstacle (VO) : ensemble des vélocités qui mènent à collision.",
 "Réciprocité : chaque agent prend la moitié de la responsabilité d'évitement.",
 "Résolution = programme linéaire de basse dimension par agent.",
 "Scalabilité : ~8 ms pour 5000 agents → temps réel massif.",
 "Adapté ici à un mouvement 1D le long d'une Lane curviligne.",
 "Comportement émergent : spirales convergentes/divergentes au centre.",
], size=13)
add_placeholder(s, Inches(8.4), Inches(2.1), Inches(4.5), Inches(3.5),
                "Schéma : demi-plan de vélocité ORCA")
add_text(s, Inches(8.4), Inches(5.75), Inches(4.5), Inches(1.3),
"Réf. : Van den Berg, Guy, Lin, Manocha, « Reciprocal n-body "
"Collision Avoidance », ISRR 2011.", size=10, color=C_DIM)
set_notes(s,
"ORCA prend une voie totalement différente. Pas de manager, pas de Claim, "
"pas de FSM : juste de la géométrie continue dans l'espace de vélocité. "
"Chaque agent calcule, pour chaque voisin, un demi-plan de vélocités "
"sécurisées. Il choisit ensuite la vélocité la plus proche de son désir, "
"contenue dans l'intersection de tous les demi-plans. La réciprocité est "
"l'axiome clé : je suppose que l'autre va aussi modifier sa vélocité, donc "
"je ne prends que la moitié de l'évitement à ma charge. Mon implémentation "
"adapte ORCA au cas 1D le long de la Lane : raisonnement en temps d'arrivée "
"plutôt qu'en vélocité libre, avec bascule cède-souple vers stop ferme si "
"le conflit devient imminent.")

# 1.8 Virtual platooning
s = slide_content("ÉTAT DE L'ART · 1.8",
                  "Virtual Platooning — Medina et al.",
                  "Projection 1D : transformer un problème 2D en suivi de leader")
add_card(s, Inches(0.6), Inches(2.1), Inches(6.1), Inches(2.7),
"Principe",
"Projeter mathématiquement position, vitesse et accélération des véhicules "
"de plusieurs voies physiques sur une voie virtuelle 1D unique traversant "
"l'intersection. Le problème devient un simple car-following.",
accent=C_ACC)
add_card(s, Inches(6.85), Inches(2.1), Inches(6.05), Inches(2.7),
"Effet visuel",
"Effet fermeture éclair (zipping) à haute vitesse, sans arrêt. Les "
"véhicules de flux orthogonaux s'entrelacent en maintenant leur cinématique. "
"Très fluide, impossible pour un conducteur humain.",
accent=C_ACC2)
add_placeholder(s, Inches(0.6), Inches(4.95), Inches(12.3), Inches(1.95),
                "Animation : effet zipping deux flux orthogonaux")
add_text(s, Inches(0.6), Inches(6.95), Inches(12.3), Inches(0.4),
"Réf. : Medina et al., « Cooperative Intersection Control Based on Virtual Platooning ».",
size=10, color=C_DIM)
set_notes(s,
"Le Virtual Platooning, c'est l'idée géniale de Medina et ses collègues. "
"On prend les véhicules de deux flux orthogonaux et on projette leur position "
"sur une voie virtuelle 1D qui traverse le carrefour. À partir de là, on "
"applique simplement de l'IDM classique : chaque véhicule suit son meneur "
"virtuel projeté. Aucun arrêt requis : si je prédis que ton arrivée projetée "
"est juste avant la mienne, tu deviens mon meneur et je cale ma vitesse "
"sur la tienne. Le résultat est un effet de fermeture éclair où deux files "
"s'entrelacent à vitesse quasi constante. C'est ce que ma VIRTUAL_PLATOON "
"policy produit, et c'est visible spectaculairement dans la simulation.")

# 1.9 MAPF
s = slide_content("ÉTAT DE L'ART · 1.9",
                  "Multi-Agent Path Finding (MAPF)",
                  "Comparaison des familles classiques")
tbl_data = [
 ["Algorithme", "Paradigme", "Résolution conflit", "Qualité"],
 ["Prioritized Planning (PP)", "Découplée séquentielle",
  "Agent low-prio évite high-prio", "Sous-optimale, rapide, sensible à l'ordre"],
 ["Priority-Based Search (PBS)", "Arborescente",
  "Explore espace des priorités partielles", "Quasi-optimale, 600 agents < 1 min"],
 ["Conflict-Based Search (CBS)", "Arbre de contraintes",
  "Contraintes spatio-temporelles", "Optimale, coût très élevé"],
 ["Token Passing (TP)", "Séquentiel par jeton",
  "Réserve chemin complet", "Incomplet, ultra-rapide (entrepôt)"],
]
left, top = Inches(0.6), Inches(2.15)
cw = [Inches(2.6), Inches(2.6), Inches(3.6), Inches(3.5)]
rh = Inches(0.75)
y = top
for r, row in enumerate(tbl_data):
    x = left
    for c, val in enumerate(row):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw[c], rh)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1B, 0x22, 0x36) if r else C_ACC
        cell.line.color.rgb = C_DIM
        cell.line.width = Pt(0.5)
        add_text(s, x + Inches(0.08), y + Inches(0.10),
                 cw[c] - Inches(0.16), rh - Inches(0.2),
                 val, size=11, bold=(r == 0),
                 color=C_BG if r == 0 else C_FG)
        x += cw[c]
    y += rh
add_text(s, Inches(0.6), Inches(6.5), Inches(12.3), Inches(0.5),
"Constat : priorités fixes ↔ deadlocks fréquents en intersection symétrique. "
"Variantes asynchrones (AWC) introduisent des priorités dynamiques par backtrack.",
size=12, color=C_DIM)
set_notes(s,
"Le Multi-Agent Path Finding est la famille théorique qui sous-tend toute "
"coordination multi-agents. PP est rapide mais sous-optimal et très sensible "
"à l'ordre. CBS est optimal mais explose en coût. PBS est un bon compromis "
"quasi-optimal. Pour mon projet, la régulation FIXED_PRIORITY correspond "
"à un Prioritized Planning trivial — axe principal toujours prioritaire. "
"Ce qui est important : la littérature démontre que les priorités fixes "
"conduisent inévitablement à des deadlocks sur les configurations symétriques. "
"C'est précisément pourquoi P2P avec sa hiérarchie de dominance et le bris "
"VIN est intéressant : il garantit la résolution déterministe de toute "
"symétrie résiduelle.")

# 1.10 MARL / theory of games
s = slide_content("ÉTAT DE L'ART · 1.10",
                  "Théorie des jeux & MARL",
                  "Frontière contemporaine : agents apprenants sous contrainte")
add_bullets(s, Inches(0.6), Inches(2.1), Inches(7.6), Inches(4.5), [
 "Intersection = jeu à somme non nulle : minimiser MON délai + éviter LA collision.",
 "Négociation = bargaining à concession monotone (convergence garantie).",
 "MARL paradigme CTDE : Centralized Training, Decentralized Execution.",
 "MACPO (IEEE 2023) : Constrained Markov Game, divergence KL bornée.",
 "Garantit politique d'apprentissage qui ne viole jamais les contraintes physiques.",
 "Performance : zéro collision, gain exponentiel vs MIP.",
 "Piste émergente : LLM-agents pour raisonnement stratégique quasi-humain.",
], size=13)
add_placeholder(s, Inches(8.4), Inches(2.1), Inches(4.5), Inches(4.8),
                "Schéma : pipeline CTDE")
set_notes(s,
"Le sommet contemporain de l'état de l'art, c'est le MARL : Multi-Agent "
"Reinforcement Learning. L'idée est de formuler le contrôle d'intersection "
"comme un jeu de Markov et d'apprendre les politiques optimales. Le paradigme "
"clé est CTDE : on entraîne avec une vision globale, on déploie de façon "
"décentralisée. Le problème historique du RL pour la sécurité critique, "
"c'est qu'un agent en apprentissage peut explorer des actions catastrophiques. "
"MACPO résout cela en bornant la divergence KL des mises à jour de politique, "
"garantissant que l'agent reste dans une région sûre. Ces approches dépassent "
"le cadre de mon implémentation actuelle, mais représentent l'horizon "
"d'extension naturel.")

# 1.11 Monte Carlo méthodologie
s = slide_content("ÉTAT DE L'ART · 1.11",
                  "Méthodologie scientifique : Monte-Carlo + métriques",
                  "Comment un simulateur devient un banc d'essai")
add_card(s, Inches(0.6), Inches(2.1), Inches(6.1), Inches(2.4),
"Hypothèse centrale",
"La performance relative entre P2P décentralisé et priorité fixe présente "
"un point d'inflexion lié à la densité du trafic. À faible densité, P2P "
"domine ; à forte densité, P2P dégénère vers four-way stop et la priorité "
"fixe redevient supérieure.",
accent=C_ACC)
add_card(s, Inches(6.85), Inches(2.1), Inches(6.05), Inches(2.4),
"Design expérimental",
"Variable indépendante : densité d'injection (0.1 → 0.8 véh/s).\n"
"Variable indépendante : type de policy (10 régulations).\n"
"N runs par point (Monte-Carlo) → moyenne + intervalle de confiance.\n"
"Exécution headless multi-thread (NullRenderer).",
accent=C_ACC2)
add_text(s, Inches(0.6), Inches(4.7), Inches(12.3), Inches(0.4),
"Métriques validées par la littérature :", size=14, bold=True, color=C_ACC)
metrics = [
 ("Throughput", "véh/min (fenêtre 60 s)"),
 ("Delay", "Δt vs free-flow"),
 ("TTC", "< 1.5 s = incident"),
 ("PET", "near-miss spatial"),
 ("Jerk", "∫|da/dt| (confort)"),
]
x = Inches(0.6)
for n, d in metrics:
    add_card(s, x, Inches(5.15), Inches(2.4), Inches(1.7), n, d, accent=C_OK)
    x += Inches(2.46)
set_notes(s,
"Pour faire de la recherche, il faut une méthodologie. La mienne est "
"directement inspirée des publications du domaine. L'hypothèse principale : "
"P2P et priorité fixe ne dominent pas dans le même régime de densité. C'est "
"une hypothèse testable. Le design expérimental balaye la densité d'injection "
"et le type de policy, on lance N runs Monte-Carlo par point pour avoir des "
"intervalles de confiance. Tout cela tourne en mode headless grâce au "
"NullRenderer : on peut faire des centaines de runs en arrière-plan sans "
"bloquer l'UI. Les métriques retenues sont les standards reconnus : débit, "
"délai, TTC pour la sécurité, PET pour les quasi-accidents, jerk pour le "
"confort passager.")

# ======================================================================
# PART 2 — DÉVELOPPEMENT
# ======================================================================
slide_section(2, "Développement & comportements",
"10 régulations · 15 fiches comportement · FSM décisionnelle")
set_notes(SLIDES[-1][0],
"Maintenant que le cadre théorique est posé, passons à ce qui est "
"concrètement implémenté dans l'application. Je vais montrer les "
"dix régulations interchangeables à chaud, puis les quinze comportements "
"de conduite individuels, et la machine à états qui les orchestre.")

# 2.1 Overview
s = slide_content("DÉVELOPPEMENT · 2.1",
                  "Vue d'ensemble de l'application",
                  "Trois états d'application + UI ImGui + rendu SFML")
add_card(s, Inches(0.6), Inches(2.1), Inches(4.0), Inches(2.0),
"MAIN_MENU",
"Sélection du scénario : catalogue prédéfini ou éditeur build mode "
"avec palette de tuiles.",
accent=C_ACC)
add_card(s, Inches(4.7), Inches(2.1), Inches(4.0), Inches(2.0),
"MONTE_CARLO",
"Banc d'essai quantitatif : balayage densité × stratégie, headless "
"multi-thread, export CSV/JSON.",
accent=C_ACC2)
add_card(s, Inches(8.8), Inches(2.1), Inches(4.1), Inches(2.0),
"SIMULATION",
"Exécution réelle 60 FPS, dashboard métriques live, overlay décisions, "
"swap policy à chaud.",
accent=C_OK)
add_placeholder(s, Inches(0.6), Inches(4.3), Inches(8.0), Inches(2.6),
                "Capture : dashboard ImGui en cours de simulation")
add_bullets(s, Inches(8.8), Inches(4.3), Inches(4.1), Inches(2.6), [
 "Cadence fixe 60 Hz",
 "FIXED_DT = 16.67 ms",
 "MAX_SUBSTEPS = 5",
 "Anti-spirale de la mort",
 "Swap policy en 1 clic",
 "Pause + diagnostic CSV",
], size=12)
set_notes(s,
"L'application est structurée autour de trois grands états. Le menu permet "
"de choisir un scénario ou de construire le sien avec une palette de tuiles. "
"Le mode Monte-Carlo lance le banc d'essai quantitatif en headless, ce qui "
"signifie que des dizaines de runs s'enchaînent sans bloquer l'UI. La "
"simulation principale tourne à 60 images par seconde avec un pas de temps "
"physique fixe de 16,67 millisecondes. Le dashboard ImGui affiche en "
"temps réel toutes les métriques et permet de swapper la régulation d'un "
"carrefour en un clic, ce qui est inestimable pour les démonstrations.")

# 2.2 — 10 régulations summary
s = slide_content("DÉVELOPPEMENT · 2.2",
                  "10 régulations implémentées",
                  "5 classiques + 5 issues de la recherche")
regs = [
 ("PRIORITY_RIGHT", "Priorité à droite, gap-acceptance directionnelle", C_ACC),
 ("STOP",          "Arrêt ≥ 0.8s puis gap-acceptance 360°",            C_ACC),
 ("YIELD",         "Cédez le passage, sans arrêt obligatoire",          C_ACC),
 ("TRAFFIC_LIGHT", "Vert/orange/rouge, dilemme orange résolu",          C_ACC),
 ("ROUNDABOUT",    "Anneau circulaire, priorité aux véhicules dedans",  C_ACC),
 ("FIXED_PRIORITY","Axe principal jamais arrêté",                       C_ACC2),
 ("P2P",           "Négociation pair-à-pair, hiérarchie de dominance",  C_ACC2),
 ("AIM",           "Réservation FCFS centralisée par carrefour",        C_ACC2),
 ("VIRTUAL_PLATOON","Insertion fermeture éclair, leader virtuel",       C_ACC2),
 ("ORCA",          "Évitement réciproque continu, mode espace ouvert",  C_ACC2),
]
cols = 5
cw = Inches(2.55)
ch = Inches(1.05)
x0 = Inches(0.4); y0 = Inches(2.15)
for i, (n, d, c) in enumerate(regs):
    r, k = divmod(i, cols)
    x = x0 + cw * k
    y = y0 + (ch + Inches(0.18)) * r
    add_card(s, x, y, cw - Inches(0.1), ch, n, d, accent=c)
add_text(s, Inches(0.4), Inches(4.85), Inches(12.6), Inches(0.5),
"Toutes implémentent IIntersectionPolicy → injectées à chaud "
"via Intersection::setRegulation(type).", size=12, color=C_DIM)
add_placeholder(s, Inches(0.4), Inches(5.45), Inches(12.6), Inches(1.5),
                "Capture : menu de swap policy dans le dashboard")
set_notes(s,
"Voici les dix régulations disponibles. Les cinq premières — priorité droite, "
"STOP, cédez, feux, rond-point — correspondent au code de la route classique. "
"Les cinq suivantes sont issues de la recherche : FIXED_PRIORITY comme ligne "
"de base de planification, P2P pour la négociation décentralisée, AIM pour la "
"réservation centralisée, Virtual Platoon pour la fermeture éclair, ORCA pour "
"l'évitement géométrique. Le point architectural important : toutes implémentent "
"la même interface IIntersectionPolicy. On peut donc swapper la régulation "
"d'un carrefour en un clic, en pleine simulation, sans recréer la géométrie. "
"C'est l'application directe du Strategy Pattern.")

# helper: behavior fiche slide
def behavior_slide(num, name, color_overlay, trigger, vit, narr):
    s = slide_content(f"COMPORTEMENT · 2.{num}",
                      f"Fiche #{num-2} — {name}",
                      f"Overlay : {color_overlay}")
    add_card(s, Inches(0.6), Inches(2.1), Inches(12.3), Inches(0.9),
             "Déclencheur", trigger, accent=C_ACC)
    add_card(s, Inches(0.6), Inches(3.15), Inches(12.3), Inches(0.9),
             "Action vitesse / leader virtuel", vit, accent=C_ACC2)
    add_card(s, Inches(0.6), Inches(4.20), Inches(12.3), Inches(1.4),
             "Narration vidéo", narr, accent=C_OK)
    add_placeholder(s, Inches(0.6), Inches(5.75), Inches(12.3), Inches(1.4),
                    "Capture / vidéo : démonstration in-app")
    set_notes(s,
        f"Le comportement {name}. Déclencheur : {trigger}. "
        f"Conséquence sur la vitesse : {vit}. "
        f"Voici la narration que je proposerais en vidéo : {narr}")
    return s

# 2.3 → 2.17 (15 fiches)
behavior_slide(3, "CRUISING (roule libre)", "⚪ blanc / transparent",
"Aucun leader détecté ET aucune intersection à freiner ET pas en virage serré ET pas en panne.",
"IDM sans leader : a = a_max · (1 − (v/v₀)^δ). v₀ = min(maxSpeed, road_limit, cornering_factor).",
"Sur autoroute libre, le véhicule accélère progressivement vers la vitesse limite, sans à-coup, en suivant fidèlement les courbes de la voie.")

behavior_slide(4, "CAR-FOLLOWING (suivi réel)", "🟠 orange",
"PerceptionResult.hasDirectObstacle = true. Filtrage par projection Frenet : |lateral| < 22 px.",
"IDM : s* = s₀ + max(0, vT + v·Δv/(2√(a·b))). Gap dynamique stable ≈ vT + s₀.",
"Le véhicule détecte le leader sur sa voie par projection curviligne, insensible au contre-sens. Il maintient un gap proportionnel à sa vitesse — signature même du modèle IDM.")

behavior_slide(5, "CORNERING (anticipation virage)", "🟡 jaune",
"Courbure locale |dθ/ds| au-dessus d'un seuil dans la fenêtre [s, s+lookahead].",
"v₀ plafonnée par v_max_curvature = √(a_lat_max / κ). Ralentissement avant, ré-accélération après.",
"Avant chaque virage, le véhicule réduit sa vitesse désirée en fonction de la courbure — comportement naturel hérité de la cinématique latérale maximale tolérée.")

behavior_slide(6, "YIELDING — Gap Acceptance", "🟣 magenta",
"Approche d'une intersection ; policy->request() retourne {shouldStop=true} ; fenêtre de conflit jugée insuffisante.",
"Leader virtuel inséré à stopLineGap devant moi, speed=0 → IDM freine vers la ligne d'arrêt.",
"À l'approche d'un carrefour à priorité à droite, le véhicule projette un véhicule fantôme à zéro km/h sur la ligne d'arrêt. L'IDM produit alors un freinage progressif identique à celui face à un vrai obstacle.")

behavior_slide(7, "STOP-HOLD (arrêt obligatoire)", "🔴 rouge clignotant",
"Intersection STOP ET (stopHeldTime < 0.8s OU currentSpeed > 5 px/s).",
"Leader virtuel à stopLineGap=0 → arrêt strict. stopHeldTime accumule, libère stopReleased=true après 0.8s à l'arrêt.",
"Au panneau STOP, le véhicule s'immobilise complètement pendant au moins huit dixièmes de seconde — comportement réglementaire — avant d'évaluer la gap-acceptance comme à un cédez le passage.")

behavior_slide(8, "TRAFFIC-LIGHT (feu + dilemme orange)", "🚦 rouge/orange",
"Intersection.getLightState(approach) ∈ {RED, ORANGE}. Tourne-à-gauche non protégé cède à l'oncoming.",
"RED → leader virtuel à la ligne. ORANGE → si d_freinage > distance restante → on PASSE ; sinon on freine.",
"Le dilemme de la zone d'orange est résolu cinématiquement : si la distance restante est inférieure à la distance de freinage confortable, le véhicule franchit prudemment plutôt que d'imposer un freinage brutal.")

behavior_slide(9, "ROUNDABOUT-INSIDE-PRIORITY", "🔵 bleu pulsant",
"Approche d'un ROUNDABOUT ; au moins un véhicule dans le disque insideDetectRadius=70 px.",
"Leader virtuel à stopLineGap jusqu'à la ligne d'entrée. Dès dans l'anneau → isCommittedToPass=true.",
"À l'approche du rond-point, le véhicule s'arrête à la ligne tant qu'un autre tourne sur l'anneau. Une fois engagé sur la couronne, il devient prioritaire et ne re-cède plus.")

behavior_slide(10, "LANE-CHANGING / OVERTAKING", "🟢 vert vif",
"Leader trop lent (speed < α·v₀) ET voie opposée libre ET créneau de rabattement libre ET personnalité agressive.",
"États NONE → OVERTAKING (lateralOffset = +largeur voie) → RETURNING → NONE. Anti-éternel via timer.",
"Lorsque le leader est trop lent et que la voie opposée est libre sur une distance suffisante, le véhicule glisse latéralement, dépasse à pleine vitesse, puis revient s'insérer dans le créneau libre devant l'ex-leader.")

behavior_slide(11, "P2P-NEGOTIATING", "🟢🔵 cyan",
"Intersection P2P ; agent en zone CLAIMING ; au moins un Claim conflictuel dominant selon la hiérarchie.",
"Si Claim dominant → cède (leader virtuel à la ligne). Sinon canEnter=true.",
"Sans aucun panneau ni feu, chaque véhicule arbitre localement : la trajectoire la plus rapide à dégager gagne ; à égalité, celle à droite ; puis la rectiligne ; puis le plus petit identifiant. L'ordre émerge sans contrôleur central.")

behavior_slide(12, "PLATOONING / VIRTUAL ZIPPING", "🟦 bleu acier",
"Intersection VIRTUAL_PLATOON ; un véhicule croisé projeté à arriver juste avant moi.",
"IDM suit un leader virtuel mobile {gap=virtualLeaderGap, speed=virtualLeaderSpeed}. Vitesse quasi constante.",
"Aucun véhicule ne s'arrête. Chacun projette ses voisins croisés sur sa propre trajectoire et se cale derrière le plus immédiat — exactement comme deux files convergeant en fermeture éclair.")

behavior_slide(13, "AIM-RESERVATION (FCFS)", "🟪 violet",
"Intersection AIM ; demande de réservation refusée (chevauchement avec slot perpendiculaire).",
"Ralentit pour repousser sa fenêtre [tEnter, tExit] et re-soumettre. reservations_[VIN] cache la demande.",
"Le carrefour agit comme une tour de contrôle : chaque véhicule réserve une fenêtre temporelle. En cas de chevauchement avec une réservation perpendiculaire, il ralentit jusqu'à obtenir un créneau libre.")

behavior_slide(14, "ORCA-SOFT-YIELD", "🟢 vert pâle",
"Intersection ORCA ; conflit non imminent ; mon temps d'arrivée > celui de l'autre + tolérance.",
"Cible = softSlowFactor · v₀ (~55%) via leader virtuel mobile. Bascule en arrêt FERME si arrivée < 0.6s.",
"Plutôt qu'un arrêt franc, le véhicule effleure sa vitesse — il se glisse dans l'interstice naturel laissé par le voisin. Le moindre risque imminent suffit à passer en arrêt ferme.")

behavior_slide(15, "KEEP-CLEAR (anti-gridlock)", "🟥 rouge bordure jaune",
"Voie de sortie de l'intersection physiquement occupée ; canEnter=true MAIS engager bloquerait.",
"Leader virtuel à la ligne d'arrêt → refus d'engager. keepClearWaited_ accumule ; bris cycle par VIN.",
"Même priorité accordée, le véhicule refuse de s'engager si l'exutoire est saturé. Cette règle 'keep clear' suffit à éviter les pâtés mexicains classiques où chacun avance d'un mètre dans la boîte.")

behavior_slide(16, "BREAKDOWN (panne)", "⚫ noir clignotant",
"Tirage RNG (forceBreakdown(seconds)) ; breakdownTimer > 0.",
"Freinage maximal → arrêt total. Les véhicules derrière le traitent comme un leader immobile.",
"En cas de panne, le véhicule devient un obstacle statique. Les véhicules suivants le détectent comme un leader à vitesse nulle et s'arrêtent à distance de sécurité — la résilience est intrinsèque au modèle de suivi.")

behavior_slide(17, "INITIALIZING (démarrage)", "🟦 bleu clair",
"currentSpeed ≈ 0 ET aucun leader devant ET voie libre derrière.",
"IDM démarre depuis 0 vers v₀, accélération maximale lissée par δ=4.",
"Au démarrage, le véhicule accélère selon la cinématique du modèle IDM — accélération maximale tant que la vitesse reste bien sous v₀, puis décroissance lisse jusqu'au régime stable.")

# 2.18 FSM hiérarchique
s = slide_content("DÉVELOPPEMENT · 2.18",
                  "FSM hiérarchique décisionnelle",
                  "ROUTING → DRIVING → NEGOTIATING → COMMITTED → AT_GOAL")
fsm = [
 ("ROUTING",      "HasPath / NoPath. Path planning A* sur la grille."),
 ("DRIVING",      "CRUISING, CAR_FOLLOWING, CORNERING, OVERTAKING, RETURNING, BREAKDOWN."),
 ("NEGOTIATING",  "ASSESS → YIELD / KEEP_CLEAR / STOP_HOLD / APPROACH_FREE."),
 ("COMMITTED",    "CROSS → EXIT. isCommittedToPass=true verrouille."),
 ("AT_GOAL",      "Destruction par la boucle principale au prochain tick."),
]
y = Inches(2.15)
for n, b in fsm:
    add_card(s, Inches(0.6), y, Inches(6.0), Inches(0.85), n, b, accent=C_ACC)
    y += Inches(0.95)
add_placeholder(s, Inches(6.8), Inches(2.15), Inches(6.1), Inches(4.6),
                "Diagramme : FSM hiérarchique mermaid")
set_notes(s,
"La FSM hiérarchique structure toute la prise de décision. Le niveau le plus "
"haut, c'est ROUTING : ai-je un chemin ? Si non, terminal. Sinon, on entre dans "
"DRIVING qui regroupe les comportements de croisière, suivi, virage, dépassement, "
"panne. Quand on approche d'une intersection, on entre dans NEGOTIATING, qui a "
"son propre sous-état ASSESS évalué à chaque tick. Une fois la ligne franchie, "
"COMMITTED verrouille via isCommittedToPass — le véhicule ne re-cède plus, "
"même si la policy change d'avis. Cette structure rend chaque transition "
"explicite et auditable.")

# 2.19 ASSESS decision tree
s = slide_content("DÉVELOPPEMENT · 2.19",
                  "Arbre de décision ASSESS",
                  "Cœur du système : où les invariants se ferment")
code = (
"ASSESS (chaque tick à l'approche) :\n"
"  ┌─ STOP & stopHeldTime < 0.8s     → STOP_HOLD\n"
"  ├─ decision = policy.request(ctx)\n"
"  │\n"
"  ├─ decision.shouldStop             → YIELD (leader virtuel)\n"
"  │\n"
"  └─ canEnter=true :\n"
"     ├─ TEST HITBOX SORTIE           → KEEP_CLEAR si insuffisant\n"
"     ├─ TEST LOOK-AHEAD ASYMÉTRIQUE  → KEEP_CLEAR si chaîne occupée\n"
"     ├─ TEST GAP-ACCEPTANCE FINAL    (Frenet, anti ghost-following)\n"
"     └─ → APPROACH_FREE\n"
"        si franchit ligne → COMMIT (isCommittedToPass=true)\n"
)
tb = add_text(s, Inches(0.6), Inches(2.1), Inches(8.4), Inches(4.6),
              code, size=12, color=C_FG, font="Consolas")
add_card(s, Inches(9.2), Inches(2.1), Inches(3.7), Inches(4.6),
"Invariants garantis",
"I1 — Aucun engagement si sortie inaccessible.\n"
"I2 — Aucun suivi de contre-sens.\n"
"I3 — Aucun freinage en plein carrefour.\n"
"I4 — Aucun freinage urgent sur orange.\n"
"I5 — Aucun cycle > KEEP_CLEAR_TIMEOUT.\n"
"I7 — Reproductibilité bit-à-bit (séquentiel).",
accent=C_OK)
set_notes(s,
"L'arbre ASSESS, c'est l'endroit où se ferment toutes les classes de bugs "
"critiques. On commence par le maintien d'arrêt pour STOP. Puis on interroge "
"la policy. Si elle dit shouldStop, on bascule en YIELD avec un leader "
"virtuel à la ligne. Si elle autorise canEnter, on ne se contente pas de "
"foncer : on teste d'abord la hitbox de la sortie pour ne pas s'engager dans "
"une boîte sans pouvoir en sortir. Puis le look-ahead asymétrique qui regarde "
"plus loin dans la chaîne, puis le gap-acceptance Frenet final. Seulement "
"après tous ces tests on déclare APPROACH_FREE. C'est cet arbre qui produit "
"les invariants listés à droite.")

# 2.20 Hitbox
s = slide_content("DÉVELOPPEMENT · 2.20",
                  "Hitbox obligatoire : anti-engagement",
                  "Aucun véhicule ne rentre s'il ne peut pas sortir")
add_bullets(s, Inches(0.6), Inches(2.1), Inches(7.0), Inches(4.2), [
 "exitLane = compute_exit_lane(self.path, intersection)",
 "requiredClear = self.bodySize.x + safetyMargin (≈ 1.5 × length)",
 "Scan occupants sur exitLane.s ∈ [0, requiredClear].",
 "Au moins un occupant → BLOCKED → KEEP_CLEAR.",
 "Sinon → engagement autorisé.",
 "Camion (length=64 px) demande 96 px ; voiture (32 px) demande 48 px.",
], size=13)
add_card(s, Inches(7.7), Inches(2.1), Inches(5.2), Inches(2.2),
"Bug évité",
"Sans cette règle : un camion s'engage dans un carrefour à priorité où la "
"sortie n'a que 50 px de libre devant le prochain véhicule → bloque "
"tous les flux pendant 10+ secondes.",
accent=C_BAD)
add_placeholder(s, Inches(7.7), Inches(4.4), Inches(5.2), Inches(2.5),
                "Capture : KEEP_CLEAR en action")
set_notes(s,
"La hitbox de sortie est obligatoire. Avant tout engagement, l'agent calcule "
"la longueur libre disponible sur la voie de sortie. S'il ne tient pas — "
"avec une marge proportionnelle à sa propre longueur — il bascule en "
"KEEP_CLEAR et reste à la ligne d'entrée. Cette règle, combinée au "
"commit-to-pass, élimine 100% des demi-engagements qui figeaient le "
"carrefour dans la version précédente. Un camion exige naturellement plus "
"de marge qu'une voiture, ce qui est exactement le comportement réaliste.")

# 2.21 Cinematic + dilemma
s = slide_content("DÉVELOPPEMENT · 2.21",
                  "Cinématique IDM & dilemme orange",
                  "d_brake = v² / (2·b_comf) — décide PASSER vs FREINER")
add_card(s, Inches(0.6), Inches(2.1), Inches(6.1), Inches(2.5),
"Dilemme zone orange",
"Si d_brake(v) > distance restante → on PASSE (freinage urgent dangereux).\n"
"Sinon → leader virtuel à la ligne, freinage confortable.\n\n"
"Évite l'écueil classique du véhicule qui freine puis re-accélère dans le "
"carrefour quand le feu passe au rouge.",
accent=C_ACC)
add_card(s, Inches(6.85), Inches(2.1), Inches(6.05), Inches(2.5),
"Override cannotStop",
"Si distance < d_brake ET decision.shouldStop :\n"
"  pendingAccel = -aMax (max brake)\n"
"  dilemmaBrakeOverride = true\n\n"
"Sans cet override, l'IDM ne voyant plus le leader virtuel "
"(déjà dans la boîte) accélérerait vers v₀.",
accent=C_ACC2)
add_placeholder(s, Inches(0.6), Inches(4.75), Inches(12.3), Inches(2.3),
                "Schéma : zone orange — décision passe/freine")
set_notes(s,
"La cinématique IDM est exploitée pour résoudre le dilemme de la zone orange. "
"On calcule la distance de freinage confortable d_brake = v carré sur 2 fois "
"la décélération confortable. Si cette distance dépasse la distance restante "
"jusqu'à la ligne, on ne peut pas s'arrêter sans freinage brutal — donc on "
"passe. Sinon, on freine doucement. L'override cannotStop gère le cas "
"limite où on est déjà trop près : on freine quand même au maximum pour "
"limiter le dépassement de la ligne, sinon l'IDM, ne voyant plus de leader "
"car déjà dans la boîte, accélérerait — comportement absurde et dangereux.")

# 2.22 Commit-to-pass
s = slide_content("DÉVELOPPEMENT · 2.22",
                  "Commit-to-pass : verrou anti freinage central",
                  "isCommittedToPass = true dès franchissement de la ligne")
add_bullets(s, Inches(0.6), Inches(2.1), Inches(7.6), Inches(4.5), [
 "Dès que position ∈ coveredTiles → isCommittedToPass = true.",
 "committedIntersectionId = id du carrefour traversé.",
 "Court-circuite toute nouvelle décision de freinage policy.",
 "Garantit qu'un changement de phase de feu en plein carrefour ne fige PAS le véhicule.",
 "Reset automatique dès sortie physique (!coveredTiles.contains(position)).",
 "Verrou par ID → multi-intersection séquentielle propre.",
], size=14)
add_card(s, Inches(8.4), Inches(2.1), Inches(4.5), Inches(4.5),
"Sans commit",
"Le feu passe à l'orange pendant que le véhicule est au centre du carrefour. "
"La policy retourne shouldStop. Sans verrou, l'agent freine au milieu, bloque "
"toutes les directions. Bug classique des simulateurs naïfs.",
accent=C_BAD)
set_notes(s,
"Le commit-to-pass est un verrou crucial. Une fois que le véhicule est "
"physiquement dans la boîte de l'intersection, on bloque toute "
"reconsidération de freinage. Sans cela, un changement de phase de feu, "
"ou un nouveau véhicule prioritaire qui apparaît, pourrait faire freiner "
"l'agent au milieu du carrefour, créant un blocage immédiat. Le reset "
"se fait automatiquement dès que l'agent sort physiquement. C'est "
"l'invariant I3.")

# 2.23 Frenet anti ghost
s = slide_content("DÉVELOPPEMENT · 2.23",
                  "Bug fix : Ghost-following éradiqué",
                  "Filtrage spatial strict par projection de Frenet")
add_bullets(s, Inches(0.6), Inches(2.1), Inches(7.6), Inches(4.5), [
 "Filtre angulaire pur (cône frontal) → faux positifs catastrophiques en virage.",
 "Solution : projection sur MA Lane curviligne (Lane::project).",
 "Filtre LATÉRAL : |proj.lateral| < 22 px (largeur voie ≈ 50 px).",
 "Filtre LONGITUDINAL : proj.s > self.s (devant).",
 "Garde-fou : |Δheading| < 45° (rejet contre-sens transitoire).",
 "Rétention du leader avec le plus petit gap.",
], size=13)
add_card(s, Inches(8.4), Inches(2.1), Inches(4.5), Inches(2.2),
"Pourquoi ça marche",
"La voie opposée a offset latéral physique > 22 px → impossible d'être retenue. "
"En virage, la projection suit la courbe. Invariant I2.",
accent=C_OK)
add_placeholder(s, Inches(8.4), Inches(4.4), Inches(4.5), Inches(2.5),
                "Schéma : projection Frenet vs cône frontal")
set_notes(s,
"Le ghost-following, c'est un véhicule qui freine pour un véhicule en "
"contre-sens. La cause racine : un filtre angulaire pur, un cône frontal, "
"qui en virage accroche la voie opposée. La solution est géométrique : on "
"projette tous les véhicules sur ma propre Lane, qui est une polyligne 1D "
"curviligne. On ne retient que ceux dont l'offset latéral est inférieur à "
"22 pixels, soit moins de la moitié d'une largeur de voie. Le résultat : "
"la voie opposée à 50 pixels est physiquement exclue. Plus aucun "
"accrochage erroné, même dans les courbes. C'est l'invariant I2.")

# 2.24 Deadlocks
s = slide_content("DÉVELOPPEMENT · 2.24",
                  "Bug fix : Deadlocks et interblocages",
                  "Look-ahead asymétrique + KEEP_CLEAR + bris min(VIN)")
add_bullets(s, Inches(0.6), Inches(2.1), Inches(7.6), Inches(4.5), [
 "Cas pathologique : 4 véhicules s'engagent, chacun bloque la sortie d'un autre.",
 "1. Look-ahead asymétrique : exige 1.5 × self.length libre devant.",
 "2. KEEP_CLEAR : reste à la ligne si exit lane saturée.",
 "3. Liveness guard : keepClearWaited_ > KEEP_CLEAR_TIMEOUT (6s).",
 "4. isCircularGridlock() : graphe orienté des attentes.",
 "5. Bris déterministe : self.vehicleId == min(VIN du cycle) → engage.",
 "VIN est un compteur global déterministe → Monte-Carlo préservé.",
], size=13)
add_card(s, Inches(8.4), Inches(2.1), Inches(4.5), Inches(4.5),
"Invariants",
"I1 — Aucun engagement sans capacité de sortie.\n\n"
"I5 — Aucun cycle ne dure plus de KEEP_CLEAR_TIMEOUT secondes.\n\n"
"Reproductibilité Monte-Carlo préservée par déterminisme du VIN.",
accent=C_OK)
set_notes(s,
"Les deadlocks sont théoriquement inévitables au-dessus d'un certain "
"seuil de densité — c'est démontré par la littérature MAPF. Pour les "
"éradiquer, j'utilise trois mécanismes combinés. D'abord le look-ahead "
"asymétrique : avant d'engager, je vérifie que la voie de sortie a "
"effectivement la place pour ma longueur. Sinon, KEEP_CLEAR : je reste "
"à la ligne même si j'ai la priorité. Si l'attente dépasse 6 secondes "
"et que la chaîne d'attentes forme un cycle, le véhicule de plus petit "
"VIN dans le cycle est désigné pour casser le deadlock en s'engageant. "
"Le VIN étant un compteur global déterministe, ce choix est identique "
"dans tous les replays — la reproductibilité Monte-Carlo est préservée.")

# 2.25 Roundabouts
s = slide_content("DÉVELOPPEMENT · 2.25",
                  "Bug fix : Ronds-points",
                  "Priorité absolue à l'anneau + commit dès la couronne")
add_bullets(s, Inches(0.6), Inches(2.1), Inches(7.6), Inches(4.5), [
 "Sans correction : véhicules sur l'anneau cèdent aux entrants → freeze.",
 "Règle 1 : tout véhicule physiquement sur l'anneau est PRIORITAIRE absolu.",
 "Règle 2 : isCommittedToPass = true dès l'entrée sur la couronne.",
 "Règle 3 : entrant cède selon véhicules INSIDE seulement.",
 "isOnRoundaboutRing(agent) = |distance(C, agent) - laneRadius| < tileSize/2.",
 "Trajectoires : arcs tangentiels + Bézier quadratique d'entrée/sortie.",
], size=13)
add_card(s, Inches(8.4), Inches(2.1), Inches(4.5), Inches(2.2),
"Invariant I9",
"Aucun véhicule sur l'anneau ne sera ralenti par un véhicule en approche.",
accent=C_OK)
add_placeholder(s, Inches(8.4), Inches(4.4), Inches(4.5), Inches(2.5),
                "Capture : rond-point à 8 véhicules fluide")
set_notes(s,
"Les ronds-points étaient le bug le plus visuellement choquant. Un véhicule "
"déjà sur l'anneau cédait à un véhicule en approche, créant une vague de "
"freinage qui figeait tout l'anneau. La correction est en trois règles. "
"Premièrement, tout véhicule physiquement sur l'anneau est prioritaire "
"absolu. Deuxièmement, on verrouille commit-to-pass dès l'entrée sur la "
"couronne. Troisièmement, l'entrant cède uniquement par rapport aux "
"véhicules déjà à l'intérieur. Le test géométrique est simple : "
"distance au centre dans la fenêtre laneRadius plus ou moins tileSize sur 2. "
"Invariant I9 : l'anneau ne se fige plus jamais.")

# ======================================================================
# PART 3 — OBSERVATIONS & EXPÉRIENCES
# ======================================================================
slide_section(3, "Observations scientifiques & expériences",
"Monte-Carlo, métriques, comportements émergents, confrontation littérature")
set_notes(SLIDES[-1][0],
"Troisième partie : les observations scientifiques. Que mesure-t-on, comment, "
"et qu'est-ce que les résultats nous disent par rapport aux prédictions de la "
"littérature.")

# 3.1 banc Monte-Carlo
s = slide_content("EXPÉRIENCES · 3.1",
                  "Banc d'essai Monte-Carlo",
                  "Du simulateur à l'instrument scientifique")
add_card(s, Inches(0.6), Inches(2.1), Inches(12.3), Inches(1.1),
"Mode headless",
"NullRenderer + thread dédié → des dizaines de runs en arrière-plan sans "
"bloquer l'UI principale. core/ 100% SFML-free.",
accent=C_ACC)
add_card(s, Inches(0.6), Inches(3.3), Inches(12.3), Inches(1.1),
"ExperimentRunner",
"Boucle : densité ∈ {0.1, 0.2, ..., 0.8} véh/s × stratégie ∈ {10 régulations} "
"× runsPerPoint. Output : vector<ResultRow>.",
accent=C_ACC2)
add_card(s, Inches(0.6), Inches(4.5), Inches(12.3), Inches(1.1),
"Mode Monte-Carlo Live",
"Injection stochastique pendant simulation visuelle. SpawnProfile pondéré "
"(wCar, wTruck, wNormal/Aggressive/Calm), gaussienne sur paramètres IDM.",
accent=C_OK)
add_placeholder(s, Inches(0.6), Inches(5.75), Inches(12.3), Inches(1.4),
                "Capture : panneau Monte-Carlo + barre progression runs")
set_notes(s,
"Le banc d'essai Monte-Carlo, c'est ce qui transforme le simulateur en "
"instrument scientifique. Il y a deux modes. Le headless : on lance des "
"dizaines de runs en arrière-plan sur un thread dédié, l'UI reste réactive, "
"et on exporte un CSV à la fin. Le live : on injecte stochastiquement des "
"véhicules dans une simulation visuelle, ce qui permet de regarder en direct "
"comment la stratégie se comporte sous charge croissante. Le SpawnProfile "
"permet de doser le mix voitures/camions et les profils de conducteurs.")

# 3.2 hypothèse
s = slide_content("EXPÉRIENCES · 3.2",
                  "Hypothèse scientifique principale",
                  "Point d'inflexion entre P2P et priorité fixe")
add_card(s, Inches(0.6), Inches(2.1), Inches(12.3), Inches(1.6),
"H1 — Inversion d'efficacité densité-dépendante",
"À faible densité, P2P décentralisé domine la priorité fixe (zéro attente "
"inutile). À forte densité (au-delà de ~0.35 véh/s/voie selon VanMiddlesworth), "
"P2P sature dans des invalidations mutuelles et dégénère vers four-way stop ; "
"la priorité fixe redevient supérieure.",
accent=C_ACC)
add_card(s, Inches(0.6), Inches(3.85), Inches(12.3), Inches(1.0),
"H2 — Métriques de sécurité",
"AIM et Virtual Platoon maintiennent un TTC > 1.5 s même à haute densité, "
"là où P2P enregistre des incidents critiques.",
accent=C_ACC2)
add_card(s, Inches(0.6), Inches(5.0), Inches(12.3), Inches(1.0),
"H3 — Confort passager",
"ORCA et Virtual Platoon minimisent le jerk accumulé vs STOP / TRAFFIC_LIGHT "
"(qui imposent des cycles arrêt/redémarrage).",
accent=C_OK)
add_text(s, Inches(0.6), Inches(6.2), Inches(12.3), Inches(0.5),
"Test statistique : différence de moyennes par densité, IC 95% sur N runs Monte-Carlo.",
size=12, color=C_DIM)
set_notes(s,
"Trois hypothèses testables. H1, c'est l'hypothèse centrale héritée de "
"VanMiddlesworth : il existe un point d'inflexion où P2P perd son avantage "
"sur les régulations à priorité fixe. H2 : les stratégies coordonnées AIM "
"et Virtual Platoon devraient garantir des marges de sécurité TTC "
"supérieures aux négociations purement locales. H3 : les stratégies "
"sans arrêt total — ORCA et Platooning — devraient produire moins de "
"jerk, donc plus de confort. Pour chaque hypothèse, on fait du test "
"statistique de différence de moyennes avec intervalle de confiance à "
"95% sur N runs Monte-Carlo.")

# 3.3 metrics implementées
s = slide_content("EXPÉRIENCES · 3.3",
                  "Métriques instrumentées",
                  "6 implémentées + 6 à étendre")
m1 = [
 ("Throughput", "véh/min (fenêtre 60s)"),
 ("Mean delay", "Δt vs free-flow"),
 ("Mean speed", "moy. véhicules actifs"),
 ("Min TTC", "min sur paires"),
 ("TTC violations", "ttc < 1.5s"),
 ("Jerk accumulé", "∫|da/dt| dt"),
]
m2 = [
 ("Block reason dist.", "% par catégorie BR"),
 ("Intersection occupancy", "% temps occupé"),
 ("Queue length", "par approche"),
 ("Mean waiting time", "par intersection"),
 ("Headway distribution", "par approche"),
 ("Deadlock incidents", "compteur bris VIN"),
]
add_text(s, Inches(0.6), Inches(2.05), Inches(6), Inches(0.4),
         "✅ Implémentées dans MetricsCollector", size=13, bold=True, color=C_OK)
add_text(s, Inches(6.85), Inches(2.05), Inches(6), Inches(0.4),
         "⬜ À ajouter (extensions)", size=13, bold=True, color=C_ACC2)
y = Inches(2.5)
for (n, d), (n2, d2) in zip(m1, m2):
    add_card(s, Inches(0.6), y, Inches(6.0), Inches(0.6), n, d, accent=C_OK)
    add_card(s, Inches(6.85), y, Inches(6.05), Inches(0.6), n2, d2, accent=C_ACC2)
    y += Inches(0.72)
set_notes(s,
"Côté implémentation, six métriques tournent déjà dans le MetricsCollector. "
"Throughput en fenêtre glissante 60 secondes, délai par rapport au flux libre, "
"vitesse moyenne, TTC minimum et compteur de violations critiques, jerk "
"accumulé pour le confort. Six autres sont des extensions documentées dans "
"Phase 5 du rapport — block reason distribution, occupancy par intersection, "
"longueur de file, temps d'attente, distribution des headways, et incidents "
"de deadlock. Toutes sont exportables en CSV et JSON pour analyse externe.")

# 3.4 Emergent behaviors
s = slide_content("EXPÉRIENCES · 3.4",
                  "Comportements émergents observés",
                  "Motifs spatio-temporels non scriptés")
emergent = [
 ("Effet accordéon dynamique",
  "Chaîne de véhicules à l'approche d'une intersection : freinage du leader → "
  "onde de freinage fluide. Si la voie se libère, ré-accélération avant arrêt "
  "complet. Démontre la string stability de l'IDM.",
  "IDM + chaîne car-following"),
 ("Effet fermeture éclair (zipping)",
  "AIM ou Virtual Platoon à haute densité : véhicules orthogonaux s'entrelacent "
  "à vitesse quasi constante. Marges d'erreur infimes — chorégraphie impossible "
  "pour conducteur humain.",
  "VIRTUAL_PLATOON / AIM"),
 ("Spirales ORCA",
  "Au centre du carrefour ORCA, les trajectoires divergent et convergent en "
  "spirales esthétiques. Émergence pure depuis la réciprocité géométrique.",
  "ORCA"),
 ("Paradoxe de Braess",
  "À haute densité P2P : invalidations mutuelles → véhicules hésitent, "
  "avancent par à-coups, finissent par mimer un four-way stop. "
  "Quantifiable via le delay.",
  "P2P saturé"),
]
y = Inches(2.1)
for n, b, tag in emergent:
    add_card(s, Inches(0.6), y, Inches(10.0), Inches(1.15), n, b, accent=C_ACC)
    add_card(s, Inches(10.7), y, Inches(2.2), Inches(1.15),
             "Trigger", tag, accent=C_OK)
    y += Inches(1.25)
set_notes(s,
"Quatre comportements émergents sont observables et reproductibles. L'effet "
"accordéon est la signature visuelle de la string stability de l'IDM : la "
"chaîne se contracte et se dilate sans onde divergente. L'effet zipping est "
"spectaculaire avec AIM ou Virtual Platoon — deux flux orthogonaux s'entrelacent "
"sans s'arrêter. Les spirales ORCA émergent purement de la géométrie réciproque, "
"sans aucune règle de fluidité programmée. Le paradoxe de Braess se manifeste "
"avec P2P à haute densité : on observe l'effondrement vers four-way stop, "
"exactement comme prédit dans la littérature.")

# 3.5 Résultats placeholder
s = slide_content("EXPÉRIENCES · 3.5",
                  "Résultats Monte-Carlo : débit vs densité",
                  "Comparaison stratégies sur carrefour isolé 4 branches")
add_placeholder(s, Inches(0.6), Inches(2.1), Inches(8.0), Inches(4.6),
                "Courbe : Throughput (véh/min) vs Densité (véh/s)\n"
                "10 stratégies superposées, IC 95%")
add_card(s, Inches(8.8), Inches(2.1), Inches(4.1), Inches(2.2),
"Lecture attendue",
"AIM domine à toutes densités.\nVirtual Platoon proche d'AIM.\n"
"P2P > FIXED en bas densité.\nP2P < FIXED en haute densité.\n"
"STOP / TRAFFIC_LIGHT baseline.",
accent=C_ACC)
add_card(s, Inches(8.8), Inches(4.4), Inches(4.1), Inches(2.3),
"Point d'inflexion",
"Croisement P2P/FIXED attendu autour de 0.35 véh/s/voie selon "
"VanMiddlesworth. À vérifier empiriquement.",
accent=C_ACC2)
set_notes(s,
"La courbe principale du banc d'essai : throughput en fonction de la densité, "
"avec les dix stratégies superposées et intervalle de confiance à 95%. La "
"lecture attendue, fondée sur la littérature : AIM domine partout grâce à sa "
"coordination optimale ; Virtual Platoon le talonne sans infrastructure "
"centralisée ; P2P est performant à basse densité mais croise FIXED_PRIORITY "
"autour de 0.35 véhicule par seconde par voie. STOP et TRAFFIC_LIGHT servent "
"de ligne de base classique. Le point d'inflexion P2P/FIXED est le test "
"empirique principal de l'hypothèse H1.")

# 3.6 Confort + sécurité
s = slide_content("EXPÉRIENCES · 3.6",
                  "Sécurité (TTC) & confort (Jerk) par stratégie",
                  "Compromis efficacité/sécurité")
add_placeholder(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(4.8),
                "Bar chart : TTC min par stratégie")
add_placeholder(s, Inches(6.9), Inches(2.1), Inches(6.0), Inches(4.8),
                "Bar chart : Jerk accumulé moyen par stratégie")
set_notes(s,
"Deux barres complémentaires. À gauche, le TTC minimum par stratégie : on "
"attend AIM, Virtual Platoon et TRAFFIC_LIGHT au-dessus du seuil 1.5 s, "
"P2P à haute densité tombant sous ce seuil. À droite, le jerk accumulé "
"moyen : ORCA et Virtual Platoon devraient minimiser le jerk car ils "
"n'imposent pas d'arrêt total ; STOP et TRAFFIC_LIGHT maximisent le jerk "
"par leurs cycles arrêt/redémarrage. Ces graphes permettent de visualiser "
"le triangle compromis débit / sécurité / confort.")

# 3.7 concordance literature
s = slide_content("EXPÉRIENCES · 3.7",
                  "Concordance avec la littérature scientifique",
                  "Phénomènes prédits vs phénomènes observés")
rows = [
 ("Effet accordéon IDM", "Treiber 2000",         "✅ Observé en chaîne car-follow", C_OK),
 ("Zipping AIM/Platoon", "Dresner 2008, Medina", "✅ Reproduit, fluide",            C_OK),
 ("Spirales ORCA",        "van den Berg 2011",    "✅ Émerge sans script",           C_OK),
 ("P2P seuil ~0.35 véh/s","VanMiddlesworth 2008", "🟡 À quantifier sur banc",        C_ACC2),
 ("Paradoxe de Braess",   "Braess 1968",          "✅ Visible à haute densité P2P", C_OK),
 ("String stability",     "Treiber 2000",         "✅ Pas d'onde divergente",        C_OK),
 ("Deadlock circulaire",  "MAPF folklore",        "✅ Reproduit puis éradiqué",     C_OK),
]
left, top = Inches(0.4), Inches(2.1)
cw = [Inches(3.6), Inches(3.0), Inches(4.8), Inches(1.4)]
rh = Inches(0.55)
hdr = ["Phénomène", "Référence", "Statut simulateur", ""]
y = top
for c, txt in enumerate(hdr):
    x = left + sum(cw[:c], Inches(0))
    cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw[c], rh)
    cell.fill.solid(); cell.fill.fore_color.rgb = C_ACC
    cell.line.fill.background()
    add_text(s, x + Inches(0.08), y + Inches(0.07), cw[c] - Inches(0.16),
             rh - Inches(0.14), txt, size=12, bold=True, color=C_BG)
y += rh
for row, c in [(r, c) for *r, c in rows]:
    x = left
    for k, val in enumerate(row):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw[k], rh)
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1B, 0x22, 0x36)
        cell.line.color.rgb = C_DIM; cell.line.width = Pt(0.5)
        col = c if k == 2 else C_FG
        add_text(s, x + Inches(0.08), y + Inches(0.07),
                 cw[k] - Inches(0.16), rh - Inches(0.14),
                 val, size=11, color=col,
                 bold=(k == 2))
        x += cw[k]
    y += rh
set_notes(s,
"Cette diapo est le test de validité scientifique. Pour chaque phénomène "
"prédit par la littérature, on vérifie son occurrence dans le simulateur. "
"L'effet accordéon, le zipping, les spirales ORCA, la string stability, "
"le deadlock circulaire et sa résolution sont tous reproduits. Le seul "
"point qui reste à quantifier rigoureusement est le seuil exact de 0.35 "
"véhicule par seconde par voie pour le décrochage P2P — c'est précisément "
"ce que le balayage Monte-Carlo va mesurer.")

# 3.8 grosse simulation issue
s = slide_content("EXPÉRIENCES · 3.8",
                  "Limitation résiduelle : grosse simulation",
                  "Des blocages aux intersections persistent à très haute densité")
add_card(s, Inches(0.6), Inches(2.1), Inches(12.3), Inches(1.7),
"Symptôme",
"Sur scénarios denses (multi-carrefours connectés, N > 200 agents), des "
"blocages localisés réapparaissent malgré KEEP_CLEAR + bris VIN. Typiquement : "
"une intersection libre, mais la voie de sortie chaîne sur un autre carrefour "
"déjà saturé en amont.",
accent=C_BAD)
add_card(s, Inches(0.6), Inches(3.9), Inches(6.1), Inches(2.0),
"Hypothèses sur la cause",
"• Look-ahead asymétrique borné à 1.5 × length → pas de propagation multi-carrefour.\n"
"• keepClearWaited_ déclenche bris VIN, mais bris LOCAL au carrefour.\n"
"• Pas de coordination inter-carrefour : approche correcte localement,\n"
"  pathologique globalement.",
accent=C_ACC2)
add_card(s, Inches(6.85), Inches(3.9), Inches(6.05), Inches(2.0),
"Pistes futures",
"• Propagation de l'occupation par graphe d'intersections.\n"
"• Hystérésis adaptative sur keepClearWaited_.\n"
"• Diagnostic CSV par-véhicule (export PAUSE) déjà utilisable pour root-cause.\n"
"• MAPF préventif sur le path globalement.",
accent=C_OK)
add_text(s, Inches(0.6), Inches(6.05), Inches(12.3), Inches(0.9),
"Sera documenté en transparence dans le rapport final. Limite honnête.",
size=12, color=C_DIM)
set_notes(s,
"Une limitation que je veux mentionner en toute transparence : sur les "
"très grosses simulations, typiquement plus de 200 agents sur un réseau "
"multi-carrefours interconnecté, des blocages localisés réapparaissent. "
"Le mécanisme local fonctionne bien sur un carrefour isolé, mais la "
"chaîne d'occupation entre plusieurs carrefours connectés produit des "
"pathologies globales que mon look-ahead local ne voit pas. La piste "
"naturelle est de propager l'information d'occupation à travers un "
"graphe d'intersections, ou d'utiliser un MAPF préventif sur les paths. "
"En attendant, le diagnostic CSV par-véhicule permet déjà d'isoler les "
"chaînes de cause pour analyse manuelle.")

# ======================================================================
# PART 4 — ARCHITECTURE
# ======================================================================
slide_section(4, "Architecture logicielle",
"Pipeline 2-phases · Strategy · diagrammes de classes · threading")
set_notes(SLIDES[-1][0],
"Quatrième partie : l'architecture logicielle. Comment le projet est "
"organisé en couches, pourquoi le pipeline 2-phases garantit l'absence de "
"course, comment le Strategy Pattern rend les régulations interchangeables, "
"et comment tout cela tient ensemble.")

# 4.1 Layers
s = slide_content("ARCHITECTURE · 4.1",
                  "Découpage en couches",
                  "core/ SFML-free ↔ render/ ↔ sim/ ↔ io/")
layers = [
 ("main.cpp",        "Orchestration : ImGui + SFML, AppState, FIXED_DT loop"),
 ("core/agent",      "IAgent, Vehicle, Car, Truck, Personality, BlockReason"),
 ("core/intersection","IIntersectionPolicy + 10 policies, Intersection oracle passif"),
 ("core/world",      "World, Lane (polyligne 1D), Tile, RoadType"),
 ("core/behavior",   "ICarFollowingModel, IdmModel, IdmParams, LeaderInfo"),
 ("core/perception", "Perception::scan() (service stateless, Frenet)"),
 ("core/metrics",    "MetricsCollector, AggregateMetrics, computeTtc"),
 ("sim",             "ThreadPool, Spawner, McLiveSession, ExperimentRunner, ParallelDecisions"),
 ("render",          "IRenderer, SfmlRenderer, NullRenderer, Camera"),
 ("io",              "ScenarioIO, SceneBuilder, ScenarioCatalog"),
]
cols = 2
cw = Inches(6.2)
ch = Inches(0.55)
x0 = Inches(0.4); y0 = Inches(2.1)
for i, (n, d) in enumerate(layers):
    r, k = divmod(i, cols)
    x = x0 + (cw + Inches(0.1)) * k
    y = y0 + (ch + Inches(0.1)) * r
    add_card(s, x, y, cw, ch, n, d, accent=C_ACC)
set_notes(s,
"Le projet est strictement découpé en couches. La règle d'or : core/ est "
"100% SFML-free. Cela signifie qu'on peut instancier tout le moteur sans "
"aucun contexte graphique, ce qui autorise le mode headless du banc "
"d'essai. Le rendu vit dans render/, l'orchestration dans sim/, la "
"persistance dans io/. Cette séparation rend chaque policy, le modèle "
"IDM, la Perception testables unitairement sans mocker SFML.")

# 4.2 Pipeline 2-phases
s = slide_content("ARCHITECTURE · 4.2",
                  "Pipeline 2-phases : compute then apply",
                  "Lecture pure puis écriture privée → safe by construction")
add_card(s, Inches(0.6), Inches(2.1), Inches(6.1), Inches(2.0),
"Phase 1 — DÉCISION",
"agent.computeDecision(agents, world) — LECTURE SEULE GLOBALE.\n"
"Produit : pendingAccel, pendingDesiredSpeed.\n"
"Parallélisable trivialement.",
accent=C_ACC)
add_card(s, Inches(6.85), Inches(2.1), Inches(6.05), Inches(2.0),
"Phase 2 — INTÉGRATION",
"agent.integrate(FIXED_DT) — ÉCRITURE PRIVÉE.\n"
"Euler symplectique : v += a·dt, s += v·dt.\n"
"Séquentielle pour bit-déterminisme.",
accent=C_ACC2)
add_text(s, Inches(0.6), Inches(4.25), Inches(12.3), Inches(0.4),
"Boucle principale 60 FPS + accumulateur + MAX_SUBSTEPS=5",
size=14, bold=True, color=C_FG)
code = (
"boucle (60 FPS) :\n"
"  pollEvents()\n"
"  simAccumulator += frameTime * simSpeedFactor\n"
"  TANT QUE simAccumulator ≥ FIXED_DT ET steps < MAX_SUBSTEPS :\n"
"    mcLive.inject()                  # Monte-Carlo spawn\n"
"    world.updateIntersections(dt)    # feux, horloges\n"
"    parallel { agent.computeDecision() }   # Phase 1\n"
"    sequential { agent.integrate(dt) }     # Phase 2\n"
"    metrics.sample(agents, dt)\n"
"    simAccumulator -= FIXED_DT\n"
"  render(world, agents, ui)\n"
)
add_text(s, Inches(0.6), Inches(4.7), Inches(12.3), Inches(2.4),
         code, size=11, color=C_FG, font="Consolas")
set_notes(s,
"Le pipeline 2-phases est l'invariant architectural fondamental. Phase 1 : "
"chaque agent prend sa décision en lisant l'état global mais sans rien "
"écrire qui sorte de ses champs pending. Phase 2 : chaque agent intègre "
"ses propres champs uniquement. Conséquence directe : la Phase 1 est "
"trivialement parallélisable car aucune course n'est possible. Tous les "
"agents décident vis-à-vis du même snapshot du monde — aucun biais d'ordre. "
"La Phase 2 reste séquentielle pour préserver le bit-déterminisme indispensable "
"aux replays Monte-Carlo. Le pas de temps fixe à 60 Hz avec accumulateur "
"garantit que les trajectoires sont identiques quel que soit le framerate.")

# 4.3 Class diagram - Vehicle
s = slide_content("ARCHITECTURE · 4.3",
                  "Diagramme de classes — Agent",
                  "IAgent, Vehicle, Car, Truck, composition")
code = (
"interface IAgent\n"
"  + computeDecision(agents, world)\n"
"  + integrate(dt)\n"
"  + getPosition() : Vec2\n"
"  + getSpeed() : float\n"
"  + getBlockReason() : BlockReason\n"
"  + getVehicleId() : int\n"
"\n"
"class Vehicle implements IAgent\n"
"  - position, velocity : Vec2\n"
"  - currentSpeed, currentAngle : float\n"
"  - currentLane : shared_ptr<Lane>\n"
"  - s : float                        // abscisse curviligne\n"
"  - pendingAccel, pendingDesiredSpeed : float\n"
"  - idm : IdmModel                   // composition\n"
"  - personality_ : Personality\n"
"  - rng_ : Rng                       // par-agent, seedé position\n"
"  - vehicleId_ : int                 // VIN déterministe\n"
"  - isCommittedToPass : bool\n"
"  - committedIntersectionId : int\n"
"  - overtakeState : OvertakeState\n"
"  - lateralOffset : float\n"
"  - currentBlockReason : BlockReason\n"
"\n"
"class Car extends Vehicle      // bodySize 32 × 16 px\n"
"class Truck extends Vehicle    // bodySize 64 × 22 px\n"
)
add_text(s, Inches(0.4), Inches(2.05), Inches(8.5), Inches(5.0),
         code, size=11, color=C_FG, font="Consolas")
add_card(s, Inches(9.1), Inches(2.05), Inches(3.8), Inches(2.4),
"Composition",
"Vehicle compose IdmModel, Personality, Rng, et référence partagée à Lane.\n"
"Pas d'héritage multiple : Strategy + composition.",
accent=C_ACC)
add_card(s, Inches(9.1), Inches(4.55), Inches(3.8), Inches(2.5),
"VIN déterministe",
"Compteur global statique → ordre de construction fixe.\n"
"Rng seedée par (startX, startY) → bit-identique run-to-run.",
accent=C_OK)
set_notes(s,
"Le diagramme de classes côté agent. IAgent est l'interface pure : "
"computeDecision, integrate, et quelques getters. Vehicle implémente "
"l'interface et compose tous les sous-systèmes : IdmModel pour le suivi, "
"Personality pour l'hétérogénéité comportementale, Rng par-agent pour "
"le déterminisme. Car et Truck héritent de Vehicle uniquement pour fixer "
"les dimensions bodySize. Aucune hiérarchie d'héritage profonde — on "
"privilégie la composition. Le VIN déterministe et la Rng seedée par "
"position garantissent que deux runs avec le même scénario produisent "
"des trajectoires bit-identiques, même si l'adresse mémoire des agents "
"change à cause d'ASLR.")

# 4.4 Class diagram - Intersection
s = slide_content("ARCHITECTURE · 4.4",
                  "Diagramme de classes — Intersection & Policies",
                  "Strategy Pattern : 10 policies interchangeables")
code = (
"class Intersection\n"
"  - id : int\n"
"  - type : RegulationType\n"
"  - coveredTiles : vector<TileCoord>\n"
"  - approaches : vector<Approach>\n"
"  - lightTimer, clock_ : float\n"
"  - policy_ : unique_ptr<IIntersectionPolicy>\n"
"  - reqMutex_ : unique_ptr<mutex>      // protège AIM\n"
"  + request(ctx) : Decision\n"
"  + setRegulation(type)                // swap hot\n"
"  + update(dt)\n"
"\n"
"interface IIntersectionPolicy\n"
"  + request(ctx, inter) : Decision\n"
"\n"
"  implementations :\n"
"    PriorityRightPolicy, StopPolicy, YieldPolicy,\n"
"    TrafficLightPolicy, RoundaboutPolicy,\n"
"    FixedPriorityPolicy, P2PPolicy, AimPolicy,\n"
"    PlatooningPolicy, OrcaPolicy\n"
)
add_text(s, Inches(0.4), Inches(2.05), Inches(8.5), Inches(5.0),
         code, size=11, color=C_FG, font="Consolas")
add_card(s, Inches(9.1), Inches(2.05), Inches(3.8), Inches(2.4),
"Oracle passif",
"L'intersection RÉPOND à request(), ne commande pas.\n"
"L'agent reste autorité finale → gap-acceptance pure.",
accent=C_ACC)
add_card(s, Inches(9.1), Inches(4.55), Inches(3.8), Inches(2.5),
"Hot swap",
"setRegulation(type) recrée la policy in-place.\n"
"Aucune destruction de géométrie, aucun reset agents.",
accent=C_OK)
set_notes(s,
"Côté intersection, le Strategy Pattern est parfaitement appliqué. "
"Intersection détient l'état géométrique : tiles couvertes, approches, "
"timer de feu. La policy_ est un unique_ptr sur l'interface "
"IIntersectionPolicy. setRegulation recrée la policy en place sans "
"toucher à la géométrie. C'est ce qui permet de swapper la régulation "
"d'un carrefour en un clic dans le dashboard. Point conceptuel important : "
"l'intersection est un oracle passif. Elle répond à request() avec une "
"structure Decision, mais ne commande jamais à l'agent. L'agent reste "
"décideur final, ce qui correspond exactement au principe de la "
"gap-acceptance en ingénierie du trafic réelle.")

# 4.5 Full class diagram placeholder
s = slide_content("ARCHITECTURE · 4.5",
                  "Vue compacte du diagramme de classes",
                  "PlantUML / Mermaid — voir docs/PHASE1_Architecture.md")
add_placeholder(s, Inches(0.4), Inches(2.1), Inches(12.6), Inches(5.0),
                "Diagramme complet : packages core::agent / core::intersection / "
                "core::world / core::perception / core::metrics / sim / render / io")
set_notes(s,
"Le diagramme complet est dans docs/PHASE1_Architecture.md, dans la section "
"1.4. Il regroupe huit packages, une vingtaine de classes principales, et "
"toutes les relations d'héritage, composition et dépendance. Pour la "
"démonstration vidéo, on peut soit l'afficher en grand, soit zoomer sur "
"des sous-vues thématiques — agent, intersection, métriques, threading.")

# 4.6 Modules & Dependencies
s = slide_content("ARCHITECTURE · 4.6",
                  "Dépendances entre modules",
                  "Flot des données et appels")
add_placeholder(s, Inches(0.4), Inches(2.1), Inches(8.4), Inches(5.0),
                "Diagramme : main → core/agent → core/perception\n"
                "                              → core/intersection → policy\n"
                "                              → core/behavior (IDM)\n"
                "                 → sim (Spawner, ExperimentRunner)\n"
                "                 → render (SfmlRenderer)\n"
                "                 → core/metrics (MetricsCollector)")
add_card(s, Inches(9.0), Inches(2.1), Inches(3.9), Inches(4.8),
"Règles",
"• core/ ne dépend de RIEN.\n"
"• render/ dépend de core (lit).\n"
"• sim/ dépend de core (orchestrer).\n"
"• io/ dépend de core (sérialise).\n"
"• main.cpp coordonne tout, sans logique métier.\n\n"
"Aucune dépendance circulaire.",
accent=C_ACC2)
set_notes(s,
"Le graphe de dépendances entre modules est acyclique. Le cœur — core/ — "
"ne dépend de rien. C'est ce qui permet de le tester unitairement sans "
"contexte graphique. render/ lit core/ pour dessiner. sim/ orchestre des "
"objets core/. io/ sérialise des objets core/. Et main.cpp coordonne tout "
"sans jamais contenir de logique métier — c'est juste un glue layer entre "
"ImGui, SFML, et le moteur. Cette discipline architecturale paye : ajouter "
"une nouvelle policy, un nouveau type d'agent, ou un nouveau format "
"d'export ne touche qu'une seule couche.")

# 4.7 Threading
s = slide_content("ARCHITECTURE · 4.7",
                  "Threading : ParallelDecisions + ThreadPool",
                  "Speedup ~8× sur Phase 1, séquentiel sur Phase 2")
add_card(s, Inches(0.6), Inches(2.1), Inches(6.1), Inches(2.3),
"Phase 1 parallèle",
"sim/ParallelDecisions::computeDecisionsParallel(agents, world, pool)\n"
"Active si N agents ≥ kParallelThreshold (150).\n"
"ThreadPool fixe, dispatch par chunks.\n"
"Speedup typique : 8× sur 8 cœurs (Phase 1 dominante).",
accent=C_ACC)
add_card(s, Inches(6.85), Inches(2.1), Inches(6.05), Inches(2.3),
"Phase 2 séquentielle",
"Conservée pour bit-déterminisme.\n"
"Coût O(n), négligeable même à N=500.\n\n"
"Monte-Carlo headless tourne sur un thread dédié, "
"sans bloquer l'UI principale.",
accent=C_ACC2)
add_card(s, Inches(0.6), Inches(4.5), Inches(12.3), Inches(2.4),
"Limitation AIM",
"AimPolicy n'est PAS bit-déterministe en multi-thread (l'ordre FCFS dépend "
"du scheduling). Documenté dans ParallelDecisions.hpp. Monte-Carlo + snapshot "
"tests tournent TOUJOURS en séquentiel. Le parallèle est réservé à "
"l'affichage temps réel des grosses scènes (N ≥ 150).",
accent=C_BAD)
set_notes(s,
"Le threading exploite la propriété fondamentale du pipeline 2-phases. "
"La Phase 1 — toutes les décisions — est parallélisée via un ThreadPool "
"fixe quand le nombre d'agents dépasse 150. On obtient un speedup quasi-"
"linéaire avec le nombre de cœurs, typiquement 8× sur 8 cœurs. La Phase 2 "
"reste séquentielle pour préserver le bit-déterminisme. La seule limitation "
"connue est AimPolicy : son ordre FCFS dépend du scheduling thread, donc "
"elle n'est pas bit-déterministe en parallèle. C'est documenté et le banc "
"Monte-Carlo tourne toujours en séquentiel pour cette raison.")

# 4.8 Reproducibility
s = slide_content("ARCHITECTURE · 4.8",
                  "Reproductibilité bit-à-bit",
                  "RNG seedée par position + VIN déterministe + Phase 2 séquentielle")
code = (
"// RNG par-agent seedé UNIQUEMENT par la position de spawn\n"
"rng_(uint64_t(startX*1000.f) * 0x9E3779B97F4A7C15ULL\n"
"   ^ uint64_t(startY*1000.f) * 0xBF58476D1CE4E5B9ULL);\n"
"\n"
"// VIN : compteur global statique, ordre de construction fixe\n"
"static int s_nextVehicleId = 0;\n"
"vehicleId_ = s_nextVehicleId++;\n"
)
add_text(s, Inches(0.4), Inches(2.1), Inches(12.6), Inches(2.0),
         code, size=12, color=C_FG, font="Consolas")
add_bullets(s, Inches(0.6), Inches(4.4), Inches(12.0), Inches(2.5), [
 "ASLR neutralisé : aucune dépendance à l'adresse mémoire des agents.",
 "Hashmap iteration order évité : pas de boucle dépendant d'un std::unordered_map dans la décision.",
 "Ordre des threads : Phase 1 parallèle, mais la Phase 2 séquentielle re-synchronise.",
 "Compteurs globaux statiques (VIN) : déterministes.",
 "→ Même scénario, même seed → trajectoires bit-identiques run-to-run.",
], size=13)
set_notes(s,
"La reproductibilité bit-à-bit est garantie par trois mécanismes combinés. "
"D'abord, chaque agent a sa propre Rng seedée déterministiquement par sa "
"position de spawn — pas d'horloge système, pas de compteur partagé. "
"Ensuite, le VIN est un compteur global statique, donc l'ordre de "
"construction est fixe et identique d'un run à l'autre. Enfin, la Phase 2 "
"séquentielle re-synchronise tous les threads. Conséquence : un même "
"scénario produit des trajectoires bit-identiques run-to-run, même avec "
"ASLR activé. C'est indispensable pour la validité statistique des "
"runs Monte-Carlo et pour les snapshot tests.")

# 4.9 Difficultés classiques contournées
s = slide_content("ARCHITECTURE · 4.9",
                  "Difficultés classiques contournées",
                  "Comment l'architecture résout les écueils SMA habituels")
rows = [
 ("O(n²) perception",      "Scope ciblé + projection Frenet bornée + Phase 1 parallèle"),
 ("Raycasting capricieux", "Projection Frenet curviligne (rejet latéral 22 px)"),
 ("Tunneling / oscillation","FIXED_DT 60 Hz + accumulateur + MAX_SUBSTEPS=5"),
 ("Deadlocks circulaires", "Look-ahead asymétrique + KEEP_CLEAR + bris min(VIN)"),
 ("Ghost-following",       "Filtre Frenet |lateral|<22 px + cap relatif < 45°"),
 ("ASLR / ordre threads",  "RNG par-position + VIN déterministe + Phase 2 séquentielle"),
 ("Couplage logique/rendu","core/ SFML-free + IRenderer + NullRenderer"),
 ("Émergent illisible",    "BlockReason exhaustif + overlay décisions + diagnostic CSV"),
]
left, top = Inches(0.4), Inches(2.15)
cw = [Inches(4.5), Inches(8.4)]
rh = Inches(0.55)
hdr = ["Difficulté", "Contournement architectural"]
y = top
for c, txt in enumerate(hdr):
    x = left + sum(cw[:c], Inches(0))
    cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw[c], rh)
    cell.fill.solid(); cell.fill.fore_color.rgb = C_ACC
    cell.line.fill.background()
    add_text(s, x + Inches(0.08), y + Inches(0.07),
             cw[c] - Inches(0.16), rh - Inches(0.14),
             txt, size=12, bold=True, color=C_BG)
y += rh
for row in rows:
    x = left
    for k, val in enumerate(row):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw[k], rh)
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1B, 0x22, 0x36)
        cell.line.color.rgb = C_DIM; cell.line.width = Pt(0.5)
        add_text(s, x + Inches(0.08), y + Inches(0.07),
                 cw[k] - Inches(0.16), rh - Inches(0.14),
                 val, size=11, color=C_FG)
        x += cw[k]
    y += rh
set_notes(s,
"Le récapitulatif des huit difficultés classiques d'une simulation "
"multi-agents de trafic, et la réponse architecturale apportée à chacune. "
"Aucune n'est contournée par hasard ou masquée par un fudge factor — "
"chacune a sa réponse explicite localisée dans une couche dédiée. C'est "
"ce qui fait la robustesse de l'ensemble.")

# ======================================================================
# PART 5 — CONCLUSION
# ======================================================================
slide_section(5, "Synthèse, limitations & travaux futurs",
"Ce qui marche, ce qui reste imparfait, où aller ensuite")
set_notes(SLIDES[-1][0],
"Dernière partie : la synthèse, la discussion honnête des limitations qui "
"subsistent, et les pistes d'extension naturelles.")

# 5.1 Synthèse
s = slide_content("CONCLUSION · 5.1",
                  "Synthèse",
                  "De la preuve de concept à la preuve de technologie")
add_bullets(s, Inches(0.6), Inches(2.1), Inches(12.3), Inches(4.5), [
 "10 régulations interchangeables, 15 comportements documentés.",
 "FSM hiérarchique stricte : 8 invariants formellement listés.",
 "Pipeline 2-phases safe-by-construction, parallélisable.",
 "Bit-déterminisme run-to-run préservé (sauf AIM en parallèle, documenté).",
 "Banc d'essai Monte-Carlo headless + live + export CSV/JSON.",
 "Métriques scientifiques standard (Throughput, Delay, TTC, PET, Jerk).",
 "Comportements émergents observés (accordéon, zipping, spirales, Braess).",
 "Trois bugs critiques éradiqués : ghost-following, deadlocks, ronds-points.",
], size=14)
set_notes(s,
"En synthèse : le simulateur couvre dix algorithmes de coordination, "
"quinze comportements de conduite, avec une machine à états formellement "
"définie et huit invariants explicites. L'architecture en pipeline "
"2-phases garantit l'absence de course par construction et permet le "
"déterminisme bit-à-bit indispensable aux runs Monte-Carlo. Tous les "
"comportements émergents prédits par la littérature ont été reproduits. "
"Trois classes de bugs critiques — ghost-following, deadlocks circulaires, "
"figeage des ronds-points — ont été éradiquées par des mécanismes "
"identifiables et localisés.")

# 5.2 Limitations
s = slide_content("CONCLUSION · 5.2",
                  "Limitations résiduelles",
                  "Discussion honnête de ce qui reste imparfait")
lim = [
 ("AIM non déterministe en multi-thread",
  "L'ordre FCFS dépend du scheduling. Monte-Carlo en séquentiel pour cette raison."),
 ("TTC en O(n²)",
  "Acceptable à N ≤ 500. Hotspot au-delà → SpatialGrid 64×64 px en suspens."),
 ("Path replanning O(n × A*)",
  "Sur édition map en build mode. Acceptable à N ≤ 100."),
 ("Pas de multi-voie réelle",
  "Une voie par direction. Dépassement par lateralOffset transitoire seulement."),
 ("Blocages multi-carrefours en grosse simulation",
  "Look-ahead local ne propage pas la chaîne d'occupation inter-carrefour."),
 ("Métriques par-intersection limitées",
  "Globales pour l'instant. Phase 5 documente les extensions (occupancy, queue, waiting)."),
]
y = Inches(2.1)
for n, b in lim:
    add_card(s, Inches(0.6), y, Inches(12.3), Inches(0.75), n, b, accent=C_BAD)
    y += Inches(0.83)
set_notes(s,
"Côté limitations : six points sont honnêtement documentés. AIM ne préserve "
"pas le bit-déterminisme en multi-thread — choix conscient pour la simplicité. "
"Le calcul TTC en O(n²) devient un hotspot au-delà de 500 agents, une grille "
"spatiale réglerait ça. Pas de gestion multi-voie réelle pour l'instant — "
"le dépassement utilise un offset latéral transitoire. Et surtout, comme "
"je l'ai mentionné, les blocages multi-carrefours en très grosse simulation "
"persistent car le look-ahead reste local. Toutes ces limitations sont "
"identifiées, localisables dans le code, et ont des pistes de résolution.")

# 5.3 Pistes futures
s = slide_content("CONCLUSION · 5.3",
                  "Travaux futurs",
                  "Extensions naturelles et frontières de recherche")
add_bullets(s, Inches(0.6), Inches(2.1), Inches(12.3), Inches(4.8), [
 "SpatialGrid 64×64 px → perception et TTC en O(n·k).",
 "Modèle MOBIL (Treiber-Hennecke-Helbing) → vrai changement de voie multi-voies.",
 "Propagation graphe inter-carrefours → résoudre les blocages globaux résiduels.",
 "Métriques par-intersection (occupancy, queue, waiting, headway) — déjà spécifiées.",
 "Apprentissage par renforcement multi-agents (MARL/MACPO) → frontière contemporaine.",
 "LLM-agents pour raisonnement stratégique quasi-humain (LLM-RL hybride).",
 "Banc d'essai étendu : import/export OpenStreetMap pour scénarios réels.",
 "Comparaison vs SUMO / Aimsun sur scénarios de référence.",
], size=14)
set_notes(s,
"Côté pistes futures, il y a un continuum d'extensions. Les plus accessibles : "
"SpatialGrid pour gagner en performance, MOBIL pour le vrai multi-voie, "
"propagation graphe pour les blocages globaux, métriques par-intersection. "
"Plus ambitieux : intégrer du MARL avec contraintes MACPO pour comparer "
"l'apprentissage automatique aux protocoles symboliques. Très ambitieux : "
"des agents pilotés par LLM pour simuler des comportements de négociation "
"quasi-humains. Et enfin, l'import de scénarios OpenStreetMap pour comparer "
"directement avec SUMO et Aimsun sur des situations réelles.")

# 5.4 References
s = slide_content("CONCLUSION · 5.4",
                  "Références bibliographiques principales",
                  "Articles fondateurs cités dans l'état de l'art")
refs = [
 "Treiber, M., Hennecke, A., Helbing, D. (2000). Congested traffic states in "
 "empirical observations and microscopic simulations. Physical Review E, 62.",
 "Dresner, K., Stone, P. (2008). A Multiagent Approach to Autonomous "
 "Intersection Management. JAIR, 31, 591-656.",
 "VanMiddlesworth, M., Dresner, K., Stone, P. (2008). Replacing the Stop Sign: "
 "Unmanaged Intersection Control for Autonomous Vehicles. AAMAS workshop.",
 "Van den Berg, J., Guy, S., Lin, M., Manocha, D. (2011). Reciprocal n-body "
 "Collision Avoidance. Robotics Research, 70, 3-19.",
 "Medina, A. et al. (2018). Cooperative Intersection Control Based on Virtual "
 "Platooning. IEEE Trans. Intelligent Transportation Systems.",
 "Stern, R. et al. (2019). Multi-Agent Path Finding: Definitions, Variants, "
 "Benchmarks. SoCS.",
 "MACPO — Multi-Agent Constrained Policy Optimization (IEEE 2023).",
 "Treiber, M., Kesting, A. (2025). Twenty-Five Years of the Intelligent "
 "Driver Model. arXiv 2506.05909.",
]
y = Inches(2.1)
for i, r in enumerate(refs):
    add_text(s, Inches(0.6), y, Inches(12.3), Inches(0.6),
             f"[{i+1}]  {r}", size=11, color=C_FG)
    y += Inches(0.55)
set_notes(s,
"Les références bibliographiques principales. La liste complète est dans le "
"document de recherche fourni dans le dossier docs, avec une cinquantaine "
"de citations académiques. Ces huit articles sont les fondamentaux qui "
"ont structuré l'implémentation. Treiber et Helbing pour l'IDM, Dresner "
"et Stone pour AIM, VanMiddlesworth pour P2P, van den Berg pour ORCA, "
"Medina pour Virtual Platooning, Stern pour MAPF, MACPO pour le MARL "
"sous contrainte, et la rétrospective des 25 ans de l'IDM pour la "
"perspective historique.")

# 5.5 Démonstration live
s = slide_content("CONCLUSION · 5.5",
                  "Démonstration live",
                  "Ce qu'on va voir tourner")
add_bullets(s, Inches(0.6), Inches(2.1), Inches(12.3), Inches(4.5), [
 "Carrefour 4 branches : swap STOP → P2P → AIM → ORCA en direct.",
 "Rond-point 8 véhicules : démonstration priorité absolue à l'anneau.",
 "Mode build : édition de map + recalculatePath() en live.",
 "Monte-Carlo live : injection stochastique, dashboard métriques à droite.",
 "Diagnostic CSV PAUSE : isoler une chaîne de blocage.",
 "Overlay décisions : couleurs comportementales en temps réel.",
], size=14)
add_placeholder(s, Inches(0.6), Inches(5.7), Inches(12.3), Inches(1.45),
                "Capture : dashboard complet + overlay décisions")
set_notes(s,
"Pour la démonstration live, je propose six scénarios. D'abord un carrefour "
"4 branches où je swap les régulations en direct pour montrer l'écart "
"comportemental : STOP, puis P2P qui négocie sans s'arrêter, puis AIM "
"avec ses créneaux, puis ORCA et son évitement continu. Ensuite le "
"rond-point à 8 véhicules pour montrer la priorité absolue à l'anneau. "
"Le mode build pour montrer la souplesse d'édition. Le Monte-Carlo "
"live pour voir les métriques se construire sous nos yeux. Et enfin "
"un diagnostic CSV sur PAUSE pour montrer comment on remonte la chaîne "
"de cause d'un blocage.")

# 5.6 Q&A
s = prs.slides.add_slide(BLANK); add_bg(s)
add_text(s, Inches(0.6), Inches(2.5), Inches(12), Inches(1.5),
         "Merci.", size=80, bold=True, color=C_FG, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.8),
         "Questions ?", size=36, color=C_ACC, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.5),
         "TrafficSimulation MAS · 2026", size=14, color=C_DIM,
         align=PP_ALIGN.CENTER)
SLIDES.append((s, "Q&A"))
set_notes(s,
"Merci de votre attention. Je suis prêt à répondre aux questions, tant "
"sur les choix scientifiques — pourquoi telle politique de coordination, "
"comment se compare à la littérature — que sur les choix architecturaux — "
"pourquoi le pipeline 2-phases, comment garantir le déterminisme, "
"comment ajouter un nouveau type d'agent ou de policy.")

# ----------------------------------------------------------------------
# Add footers (page numbers)
# ----------------------------------------------------------------------
total = len(SLIDES)
for i, (sl, chap) in enumerate(SLIDES):
    add_footer(sl, i + 1, total, chap)

# ----------------------------------------------------------------------
# Save
# ----------------------------------------------------------------------
out = r"C:\Users\fddav\Desktop\TrafficSimulation\TrafficSimulation_Presentation.pptx"
prs.save(out)
print("SAVED:", out)
print("Slides:", total)
